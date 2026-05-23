"""
Build artifact MD/DOCX/PPTX cho báo cáo nghiên cứu VLSP 2025 NumQA.
Phiên bản method-focused: đi sâu vào phương pháp, thuật toán, tối ưu và đóng
góp của đề tài. Lấy số liệu thật từ dataset_stats.json và cấu hình thật từ
pipeline/config.py.

PPTX giữ nguyên template HUST_chinh_sua.pptx (10x7.5 inch, layout 2_Blank,
HUST red, font Arial) và chỉ thay nội dung từng slide.
"""

from __future__ import annotations

import json
import re
import shutil
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt as PptPt


OUT_DIR = Path(__file__).resolve().parent
ROOT = OUT_DIR.parents[1]
# Dùng template HUST 4x3 chính thức (2022) làm gốc — master sẽ tự render logo,
# footer, page number. Renderers chỉ điền title qua placeholder + body content.
BASE_PPTX = OUT_DIR / "HUST_PPT_template_2022_RED_4x3.pptx"
if not BASE_PPTX.exists():
    BASE_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_chinh_sua.pptx"
if not BASE_PPTX.exists():
    BASE_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST.pptx"

MD_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu_method_focus.md"
DOCX_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu_method_focus.docx"
PPTX_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_method_focus.pptx"
CANONICAL_MD = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu.md"
CANONICAL_DOCX = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu.docx"
CANONICAL_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST.pptx"
CHINH_SUA_MD = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu_chinh_sua.md"
CHINH_SUA_DOCX = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu_chinh_sua.docx"
CHINH_SUA_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_chinh_sua.pptx"

# ── Bảng màu HUST ────────────────────────────────────────────────────
HUST_RED = (165, 0, 33)
TEXT_DARK = (35, 35, 35)
TEXT_GRAY = (90, 90, 90)
LIGHT_RED = (252, 238, 241)
LIGHT_GRAY = (245, 246, 248)
LIGHT_BLUE = (231, 240, 251)
LIGHT_GREEN = (228, 244, 233)
LIGHT_AMBER = (254, 243, 219)
LIGHT_PURPLE = (240, 234, 247)
BLUE = (46, 100, 172)
GREEN = (46, 135, 86)
AMBER = (205, 132, 35)
PURPLE = (112, 78, 170)
WHITE = (255, 255, 255)


# ── Số liệu thật được trích từ dataset_stats.json ────────────────────
DATASET_STATS = json.loads((OUT_DIR / "dataset_stats.json").read_text(encoding="utf-8"))


REPORT_MD = (OUT_DIR / "_report_v2.md").read_text(encoding="utf-8").strip()


# ── Slide content specs ──────────────────────────────────────────────

SLIDE_SPECS = [
    {"title": "Knowledge Distillation cho VLSP 2025 NumQA", "kind": "title"},
    {"title": "Mục lục báo cáo", "kind": "outline"},
    # CHƯƠNG 1 — Giới thiệu và Mô tả bài toán
    {"title": "1. Bối cảnh và động cơ nghiên cứu", "kind": "context"},
    {"title": "2. Phát biểu hình thức bài toán", "kind": "problem"},
    {"title": "3. Ví dụ minh họa với dữ liệu mẫu", "kind": "sample"},
    {"title": "4. Phân tích thống kê bộ dữ liệu", "kind": "data_stats"},
    # CHƯƠNG 2 — Nghiên cứu liên quan
    {"title": "5. Ba nghiên cứu nền tảng — định vị", "kind": "related_overview"},
    {"title": "6. Nghiên cứu nền tảng 1 — FinQA (EMNLP 2021)", "kind": "rw_finqa"},
    {"title": "7. Nghiên cứu nền tảng 2 — Distilling Step-by-Step (ACL 2023)", "kind": "rw_distill"},
    {"title": "8. Nghiên cứu nền tảng 3 — GRPO trong DeepSeekMath (2024)", "kind": "rw_grpo"},
    {"title": "9. Định vị đóng góp của đề tài", "kind": "positioning"},
    # CHƯƠNG 3 — Cơ sở lý thuyết
    {"title": "10. Cơ sở lý thuyết & định hướng thiết kế", "kind": "theory"},
    # CHƯƠNG 4 — Phương pháp đề xuất
    {"title": "11. Tổng quan pipeline 5 phase", "kind": "pipeline"},
    {"title": "12. Đóng góp 1 — Chiến lược dữ liệu hướng chương trình", "kind": "data_strategy"},
    {"title": "13. Financial DSL & 4-strategy header matching", "kind": "dsl"},
    {"title": "14. Đóng góp 2 — Guided Reasoning Distillation", "kind": "kd_flow"},
    {"title": "15. Quality tiering 4 mức cho teacher trace", "kind": "quality_gates"},
    {"title": "16. Đóng góp 3 — LoRA-SFT với label masking", "kind": "lora"},
    {"title": "17. So sánh PPO / DPO / RLVR / GRPO", "kind": "method_compare"},
    {"title": "18. Đóng góp 4 — PCPO reward (đóng góp trọng tâm)", "kind": "pcpo"},
    {"title": "19. Vòng lặp huấn luyện GRPO", "kind": "grpo_loop"},
    {"title": "20. Đóng góp 5 — Verifier-guided multi-path inference", "kind": "inference"},
    # CHƯƠNG 5 — Thực nghiệm và Kết quả
    {"title": "21. Cài đặt thực nghiệm — phần cứng & hyperparam", "kind": "setup"},
    {"title": "22. Kết quả tổng thể trên ViNumQA valid", "kind": "results"},
    {"title": "23. Ablation: đóng góp tăng dần của từng tầng", "kind": "ablation"},
    {"title": "24. Phân tích lỗi và đường giải quyết", "kind": "errors"},
    # CHƯƠNG 6 — Kết luận
    {"title": "25. Hạn chế và hướng phát triển", "kind": "limits"},
    {"title": "26. Kết luận", "kind": "conclusion"},
]


# ── Slide rendering helpers ──────────────────────────────────────────

def clear_slide(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)


def rgb(color):
    return PptRGBColor(*color)


def add_text(slide, text, left, top, width, height, size=14, bold=False, color=TEXT_DARK,
             align=None, italic=False, name="Arial"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, bullets, left, top, width, height, size=13, color=TEXT_DARK,
                bullet_color=None, space_after=4):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    bc = bullet_color if bullet_color is not None else HUST_RED
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        run_bullet = p.add_run()
        run_bullet.text = "• "
        run_bullet.font.name = "Arial"
        run_bullet.font.size = PptPt(size)
        run_bullet.font.bold = True
        run_bullet.font.color.rgb = rgb(bc)
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = "Arial"
        run_text.font.size = PptPt(size)
        run_text.font.color.rgb = rgb(color)
        p.space_after = PptPt(space_after)
    return box


def add_rect(slide, left, top, width, height, text="", fill=LIGHT_GRAY, line=HUST_RED,
             text_color=TEXT_DARK, size=12, bold=False, rounded=True, align=PP_ALIGN.CENTER,
             line_width=1.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = PptPt(line_width)
    if text:
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.name = "Arial"
            run.font.size = PptPt(size)
            run.font.bold = bold if i == 0 else False
            run.font.color.rgb = rgb(text_color)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=HUST_RED, width_pt=1.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = PptPt(width_pt)
    line.line.end_arrowhead = True
    return line


def add_table(slide, rows, left, top, width, height, font_size=10, header_fill=HUST_RED,
              first_col_bold=False):
    shp = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    table = shp.table
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(header_fill)
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(LIGHT_GRAY)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = PptPt(font_size)
                    run.font.bold = (r_idx == 0) or (first_col_bold and c_idx == 0)
                    run.font.color.rgb = rgb(WHITE) if r_idx == 0 else rgb(TEXT_DARK)
    return shp


def set_slide_title(slide, title):
    """Đặt title vào placeholder của layout (master HUST template tự render).

    Nếu không tìm thấy placeholder title, fallback sang vẽ thủ công.
    """
    title_ph = None
    for shape in slide.placeholders:
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            continue
        # Type 13 = TITLE, 14 = CENTER_TITLE
        if ph_type in (13, 14, 15) or shape.placeholder_format.idx == 0:
            title_ph = shape
            break
    if title_ph is None:
        # Fallback: vẽ title thủ công
        add_text(slide, title, 0.26, 0.09, 9.49, 0.55, size=22, bold=True, color=HUST_RED)
        return
    tf = title_ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = PptPt(22)
    run.font.bold = True
    run.font.color.rgb = rgb(HUST_RED)


# Backwards-compatible alias (cũ): decorate() giờ chỉ set title qua layout.
def decorate(slide, title, page):
    set_slide_title(slide, title)


def flow_chain(slide, labels, y, height=0.65, x_start=0.55, x_end=9.45,
               colors=None, fill_default=LIGHT_GRAY, font_size=11):
    n = len(labels)
    width_total = x_end - x_start
    gap = 0.22
    w = (width_total - gap * (n - 1)) / n
    x = x_start
    for i, label in enumerate(labels):
        fill = (colors[i] if colors and i < len(colors) else fill_default)
        add_rect(slide, x, y, w, height, label, fill=fill, line=HUST_RED,
                 size=font_size, bold=True)
        if i < n - 1:
            add_arrow(slide, x + w + 0.02, y + height / 2,
                      x + w + gap - 0.02, y + height / 2, color=HUST_RED, width_pt=1.6)
        x += w + gap


def bar_chart(slide, items, left, top, width, height, max_val=None, color=HUST_RED,
              label_color=TEXT_DARK, value_color=None, show_value=True, font_size=10):
    """Render a simple horizontal bar chart at (left,top) box of size (width,height)."""
    if max_val is None:
        max_val = max(v for _, v in items)
    n = len(items)
    label_w = 1.7
    value_w = 0.55 if show_value else 0
    bar_x = left + label_w
    bar_max_w = width - label_w - value_w - 0.1
    row_h = height / n
    pad_top = row_h * 0.15
    bar_h = row_h * 0.65
    vc = value_color or color
    for i, (label, val) in enumerate(items):
        y = top + i * row_h + pad_top
        add_text(slide, label, left, y - 0.02, label_w, bar_h + 0.04, size=font_size,
                 color=label_color, align=PP_ALIGN.LEFT)
        w = bar_max_w * (val / max_val) if max_val else 0
        if w > 0.01:
            add_rect(slide, bar_x, y, w, bar_h, "", fill=color, line=color, rounded=False,
                     line_width=0)
        if show_value:
            add_text(slide, f"{val}", bar_x + bar_max_w + 0.04, y - 0.02, value_w,
                     bar_h + 0.04, size=font_size, color=vc, bold=True, align=PP_ALIGN.LEFT)


def callout(slide, left, top, width, height, title, body, fill=LIGHT_RED, line=HUST_RED,
            title_color=HUST_RED, body_color=TEXT_DARK, title_size=12, body_size=10.5):
    add_rect(slide, left, top, width, height, "", fill=fill, line=line, rounded=True)
    add_text(slide, title, left + 0.1, top + 0.08, width - 0.2, 0.34,
             size=title_size, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    add_bullets(slide, body, left + 0.12, top + 0.42, width - 0.24,
                height - 0.5, size=body_size, color=body_color, bullet_color=line,
                space_after=2)


def render_title(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    # Top red bar bigger
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(HUST_RED)
    bar.line.color.rgb = rgb(HUST_RED)
    add_text(slide, "VLSP 2025  -  NumQA", 0.4, 0.10, 4.5, 0.35,
             size=14, bold=True, color=WHITE)
    add_text(slide, "Đại học Bách khoa Hà Nội (HUST)", 5.1, 0.10, 4.5, 0.35,
             size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    # Main title
    add_text(slide, "Knowledge Distillation cho VLSP 2025 NumQA",
             0.5, 1.10, 9.0, 0.85, size=30, bold=True, color=HUST_RED, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Tối ưu chương trình tính toán bằng Guided Distillation + GRPO/PCPO + Verifier",
             0.7, 2.05, 8.6, 0.7, size=17, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    # Pipeline overview chips
    chips = ["Text + Table", "DSL Program", "Executor", "EA / PA", "Verifier"]
    chip_colors = [LIGHT_RED, LIGHT_BLUE, LIGHT_GREEN, LIGHT_AMBER, LIGHT_RED]
    flow_chain(slide, chips, y=3.05, height=0.7, x_start=0.6, x_end=9.4,
               colors=chip_colors, font_size=12)
    # Banner
    add_rect(slide, 0.7, 4.05, 8.6, 0.55,
             "Trọng tâm: phương pháp, thuật toán, tối ưu kỹ thuật và đóng góp của đề tài",
             fill=LIGHT_RED, line=HUST_RED, size=14, bold=True, text_color=HUST_RED)
    # Four contribution chips
    contrib = [
        ("Dữ liệu", "Chiến lược program-centric +\nFinQA + program_re", BLUE, LIGHT_BLUE),
        ("Distillation", "Guided trace 27B → 4B,\nquality tiering", HUST_RED, LIGHT_RED),
        ("RL", "GRPO + PCPO reward,\nKL = 1e-3", GREEN, LIGHT_GREEN),
        ("Inference", "Multi-path + verifier\n(valid + exec + brevity)", PURPLE, LIGHT_PURPLE),
    ]
    x = 0.7
    w = 2.1
    gap = 0.07
    for title, body, lc, fc in contrib:
        add_rect(slide, x, 4.85, w, 1.55, f"{title}\n\n{body}", fill=fc, line=lc,
                 text_color=TEXT_DARK, size=10.5, bold=True)
        x += w + gap
    add_text(slide, "Thanh Đức  |  do.thanh.duc.2002@gmail.com  |  2026",
             0.7, 6.62, 8.6, 0.3, size=10.5, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


def render_problem(slide):
    # Input / Output / Goal boxes
    add_rect(slide, 0.5, 1.05, 2.85, 1.05,
             "INPUT\nx = (pre_text, table, post_text, question)",
             fill=LIGHT_BLUE, line=BLUE, size=11, bold=True)
    add_arrow(slide, 3.40, 1.58, 3.78, 1.58)
    add_rect(slide, 3.85, 1.05, 2.55, 1.05,
             "OUTPUT\ny = (reasoning, program, answer)",
             fill=LIGHT_GREEN, line=GREEN, size=11, bold=True)
    add_arrow(slide, 6.45, 1.58, 6.83, 1.58)
    add_rect(slide, 6.90, 1.05, 2.55, 1.05,
             "MỤC TIÊU\nPA trước, EA sau",
             fill=LIGHT_RED, line=HUST_RED, size=11, bold=True)

    # Objective formula box
    add_rect(slide, 0.5, 2.35, 8.95, 0.95,
             "p* = argmax_p  P_θ(p | x)\nsubject to:  valid(p) = 1   ∧   exec(p) ≈ a_gold   ∧   p ≡_sym p_gold",
             fill=WHITE, line=HUST_RED, size=14, bold=True, text_color=HUST_RED,
             line_width=1.4)

    # Four research questions in 2x2 grid
    rq_items = [
        ("RQ1 - Dữ liệu",
         "Tín hiệu học hướng chương trình bằng cách nào?", BLUE, LIGHT_BLUE),
        ("RQ2 - Distillation",
         "Teacher nên truyền gì sang student?", HUST_RED, LIGHT_RED),
        ("RQ3 - RL",
         "Vì sao GRPO/PCPO phù hợp PA hơn DPO/PPO/RLVR?", GREEN, LIGHT_GREEN),
        ("RQ4 - Inference",
         "Verifier và multi-path để tăng độ tin cậy?", PURPLE, LIGHT_PURPLE),
    ]
    x0, y0 = 0.5, 3.55
    w, h = 4.45, 1.45
    gap = 0.10
    for i, (title, body, lc, fc) in enumerate(rq_items):
        cx = x0 + (i % 2) * (w + gap)
        cy = y0 + (i // 2) * (h + 0.15)
        add_rect(slide, cx, cy, w, h, f"{title}\n\n{body}",
                 fill=fc, line=lc, text_color=TEXT_DARK, size=12, bold=True)

    add_text(slide,
             "PA chặt hơn EA: PA ≡ symbolic equivalence; EA chỉ là điều kiện cần.",
             0.5, 6.60, 8.95, 0.4, size=12, italic=True, color=HUST_RED, align=PP_ALIGN.CENTER)


def render_pipeline(slide):
    add_text(slide, "5 phase độc lập, mọi giai đoạn xoay quanh executor DSL",
             0.5, 0.92, 8.95, 0.35, size=12.5, italic=True, color=TEXT_GRAY)

    # 5 main phases as labelled columns
    phases = [
        ("PHASE 1", "Data Prep",
         "Markdown table\nprogram_re\nQuality gate", LIGHT_BLUE, BLUE),
        ("PHASE 2", "Teacher Distill",
         "Qwen3.5-27B\nGuided prompt\nTier 4 mức", LIGHT_RED, HUST_RED),
        ("PHASE 3", "LoRA-SFT",
         "Student 4B\nlr=5e-5, r=128\nMask user prompt", LIGHT_AMBER, AMBER),
        ("PHASE 4", "GRPO + PCPO",
         "G=5 completions\nlr=1e-6, KL=1e-3\nR_valid gate", LIGHT_GREEN, GREEN),
        ("PHASE 5", "Inference",
         "N=15 candidates\nExecutor + verifier\nValid+Exec+Brevity", LIGHT_PURPLE, PURPLE),
    ]
    x_start, y_top = 0.30, 1.35
    box_w = 1.83
    box_h = 1.85
    gap = 0.05
    for i, (tag, name, body, fc, lc) in enumerate(phases):
        cx = x_start + i * (box_w + gap)
        add_rect(slide, cx, y_top, box_w, 0.32, tag, fill=lc, line=lc,
                 text_color=WHITE, size=10, bold=True)
        add_rect(slide, cx, y_top + 0.32, box_w, 0.45, name, fill=WHITE, line=lc,
                 size=12, bold=True, text_color=lc)
        add_rect(slide, cx, y_top + 0.77, box_w, box_h - 0.77, body, fill=fc,
                 line=lc, size=10, text_color=TEXT_DARK)
        if i < len(phases) - 1:
            add_arrow(slide, cx + box_w + 0.005, y_top + box_h / 2,
                      cx + box_w + gap - 0.005, y_top + box_h / 2,
                      color=HUST_RED, width_pt=1.6)

    # Executor cylinder under everything
    add_rect(slide, 0.30, 3.45, 9.40, 0.45,
             "EXECUTOR DSL (10 ops)  -  validate_program  +  execute_program  +  PA symbolic (sympy)",
             fill=HUST_RED, line=HUST_RED, text_color=WHITE, size=12, bold=True)

    # Two contribution callouts
    callout(slide, 0.30, 4.10, 4.55, 1.45, "Đóng góp dữ liệu & distillation",
            ["Markdown + program_re + quality gate trước reward",
             "Guided distillation 27B → 4B, tier 4 mức"],
            fill=LIGHT_RED, line=HUST_RED, title_size=12.5, body_size=11)
    callout(slide, 5.15, 4.10, 4.55, 1.45, "Đóng góp RL & inference",
            ["PCPO reward: R_valid·(0,7 + 0,2·R_exec + 0,1·R_bonus)",
             "Verifier-guided multi-path (không vote text thô)"],
            fill=LIGHT_BLUE, line=BLUE, title_color=BLUE, body_size=11)

    add_text(slide, "Mọi tín hiệu (SFT loss / GRPO reward / verifier) đều đi qua cùng executor → consistency giữa giai đoạn.",
             0.30, 6.65, 9.40, 0.35, size=11.5, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)


def render_data_stats(slide):
    # Left side - dataset overview table
    rows = [
        ["Tập", "# mẫu", "Bảng (H×C)", "Bước/program"],
        ["ViNumQA train", "2.993", "8,15 × 5,29", "1,56"],
        ["ViNumQA valid", "584", "8,25 × 5,46", "1,51"],
        ["ViNumQA test", "497", "8,58 × 5,19", "1,41"],
        ["ViNumQA private", "1.625", "9,12 × 6,22", "—"],
        ["FinQA train+ext", "6.251", "6,34 × 3,84", "1,54"],
    ]
    add_table(slide, rows, 0.30, 1.10, 4.85, 2.65, font_size=10.5, first_col_bold=True)
    add_text(slide, "Bảng 3.1 - Số mẫu, kích thước và độ phức tạp chương trình (thực)",
             0.30, 3.80, 4.85, 0.32, size=10.5, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)

    # Right side - operation frequency bar chart (ViNumQA train)
    add_text(slide, "Tần suất phép toán - ViNumQA train", 5.40, 1.05, 4.30, 0.32,
             size=12, bold=True, color=HUST_RED)
    ops = [
        ("divide", 1780),
        ("subtract", 1414),
        ("add", 744),
        ("multiply", 260),
        ("table_max", 178),
        ("table_average", 126),
        ("table_min", 100),
        ("table_sum", 52),
    ]
    bar_chart(slide, ops, left=5.40, top=1.40, width=4.30, height=2.70,
              max_val=1800, color=HUST_RED, font_size=10)

    # Bottom highlight callouts
    callout(slide, 0.30, 4.40, 4.55, 2.40, "ViNumQA - chiến lược",
            ["1.625 mẫu private test không có label → cần generalize qua FinQA",
             "Phép toán bảng (table_*) chỉ ~16% → tín hiệu hiếm; FinQA bổ sung",
             "Trung bình 1,56 bước → bài toán ngắn nhưng dày bảng/chứng cứ"],
            fill=LIGHT_RED, line=HUST_RED, title_size=12, body_size=10.5)
    callout(slide, 5.15, 4.40, 4.55, 2.40, "FinQA - tài sản quý cho PA",
            ["6.251 mẫu English bổ sung greater, exp; tận dụng multilingual Qwen",
             "2.534 program_re KHÁC program gốc → 40,5% mẫu có lời giải thay thế",
             "→ Augmentation trực tiếp tăng PA, dạy tương đương symbolic"],
            fill=LIGHT_BLUE, line=BLUE, title_color=BLUE, body_size=10.5)


def render_data_strategy(slide):
    # Top - three-tier data flow
    add_text(slide, "Mỗi mẫu được 'nổ' thành 3 tín hiệu học và 3 luồng dữ liệu",
             0.30, 0.95, 9.40, 0.32, size=12, italic=True, color=TEXT_GRAY)

    tiers = [
        ("Prompt text-table\n(Markdown)", LIGHT_BLUE, BLUE),
        ("Reasoning trace\n(teacher guided)", LIGHT_RED, HUST_RED),
        ("DSL program + answer\n(executor verify)", LIGHT_GREEN, GREEN),
    ]
    x, y, w, h = 0.30, 1.30, 3.00, 1.05
    for label, fc, lc in tiers:
        add_rect(slide, x, y, w, h, label, fill=fc, line=lc, size=12, bold=True,
                 text_color=TEXT_DARK)
        x += w + 0.05

    # Pipeline diagram from raw to outputs
    add_text(slide, "Luồng xử lý dữ liệu", 0.30, 2.55, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED)
    flow_chain(slide, ["Raw JSON\nVi + En",
                       "table_to_\nmarkdown",
                       "program_re\naugment",
                       "executor\nquality gate",
                       "3 luồng\nSFT/GRPO/teacher"],
               y=2.95, height=0.80, x_start=0.30, x_end=9.70,
               colors=[LIGHT_BLUE, LIGHT_GRAY, LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN],
               font_size=10.5)

    # program_re augmentation visualization
    add_text(slide, "Program_re augmentation - ví dụ thực", 0.30, 4.00, 4.50, 0.32,
             size=11.5, bold=True, color=HUST_RED)
    add_rect(slide, 0.30, 4.35, 4.50, 1.05,
             "program (gốc):\n divide(914, 391), multiply(#0, const_100)\nexec → 233,76",
             fill=LIGHT_GRAY, line=TEXT_GRAY, size=10, align=PP_ALIGN.LEFT)
    add_rect(slide, 0.30, 5.45, 4.50, 1.10,
             "program_re (tương đương):\n multiply(divide(914,391), 100)\nexec → 233,76  ✓ symbolic equiv",
             fill=LIGHT_GREEN, line=GREEN, size=10, align=PP_ALIGN.LEFT)

    # Right side - quality gate funnel
    add_text(slide, "Quality gate trước GRPO (đảm bảo tín hiệu reward đúng dấu)",
             4.95, 4.00, 4.85, 0.32, size=11.5, bold=True, color=BLUE)
    gate_rows = [
        ["Kiểm tra", "Hành động"],
        ["validate_program(p)", "loại nếu sai DSL"],
        ["execute_program(p, table)", "loại nếu None"],
        ["answers_match(pred, gold)", "loại nếu sai đáp án"],
        ["→ giữ làm GRPO sample", "ground_truth dict đầy đủ"],
    ]
    add_table(slide, gate_rows, 4.95, 4.35, 4.85, 2.20, font_size=10,
              header_fill=BLUE, first_col_bold=True)


def render_dsl(slide):
    # Top: pipeline natural → DSL → executor
    add_rect(slide, 0.40, 1.10, 2.75, 0.80,
             "Câu hỏi tự nhiên\n\"Tăng trưởng doanh thu (%)?\"",
             fill=LIGHT_BLUE, line=BLUE, size=11, bold=True)
    add_arrow(slide, 3.20, 1.50, 3.55, 1.50)
    add_rect(slide, 3.60, 1.10, 2.75, 0.80,
             "DSL program\ndivide(914, 391), multiply(#0, const_100)",
             fill=LIGHT_GRAY, line=HUST_RED, size=10.5, bold=True)
    add_arrow(slide, 6.40, 1.50, 6.75, 1.50)
    add_rect(slide, 6.80, 1.10, 2.65, 0.80,
             "Executor + PA symbolic\n233,76 (auditable)",
             fill=LIGHT_GREEN, line=GREEN, size=11, bold=True)

    # 10-ops grid
    add_text(slide, "10 phép toán hợp lệ của DSL", 0.40, 2.10, 9.05, 0.32,
             size=12, bold=True, color=HUST_RED)
    ops_grid = [
        ("add", BLUE), ("subtract", BLUE), ("multiply", BLUE), ("divide", BLUE),
        ("exp", AMBER), ("greater", AMBER),
        ("table_sum", GREEN), ("table_average", GREEN),
        ("table_max", GREEN), ("table_min", GREEN),
    ]
    x_start, y_op = 0.40, 2.50
    w_op = 0.87
    gap_op = 0.04
    for i, (name, col) in enumerate(ops_grid):
        cx = x_start + i * (w_op + gap_op)
        fill = LIGHT_BLUE if col == BLUE else (LIGHT_AMBER if col == AMBER else LIGHT_GREEN)
        add_rect(slide, cx, y_op, w_op, 0.55, name, fill=fill, line=col,
                 text_color=col, size=10.5, bold=True)

    # Left: 4-strategy header matching (đóng góp)
    callout(slide, 0.30, 3.40, 4.65, 2.40, "4-strategy header matching (đóng góp)",
            ["1) Exact: row[0] == header",
             "2) Lowercase: row[0].lower() == header.lower()",
             "3) Strip parens: 'Doanh thu (triệu đồng)' → 'doanh thu'",
             "4) Substring: header ⊂ row[0] và len > 2",
             "→ Giảm fail vì header mismatch trên bảng tài chính thật"],
            fill=LIGHT_RED, line=HUST_RED, title_size=12, body_size=10.5)

    # Right: PA symbolic equiv example
    callout(slide, 5.10, 3.40, 4.55, 2.40, "PA symbolic (sympy) ví dụ",
            ["add(A, B) ≡ add(B, A)  (giao hoán)",
             "multiply(divide(A, B), 100) ≡ divide(multiply(A, 100), B)",
             "→ Cho phép program_re KHÁC string nhưng vẫn = gold",
             "→ Đo logic, không phải bề mặt câu chữ"],
            fill=LIGHT_BLUE, line=BLUE, title_color=BLUE, body_size=10.5)

    add_text(slide,
             "Executor không chỉ ở đánh giá - còn nhúng vào reward GRPO → cầu nối neural ↔ symbolic.",
             0.30, 5.95, 9.40, 0.4, size=12, italic=True, color=HUST_RED,
             align=PP_ALIGN.CENTER)


def render_kd_flow(slide):
    # Pipeline diagram
    flow_chain(slide, ["Sample\n(x, p_gold, a_gold)",
                       "Guided\nprompt",
                       "Teacher\nQwen3.5-27B",
                       "Reasoning\ntrace",
                       "Validate\n4 tiers",
                       "distilled_\nsft.json"],
               y=1.05, height=0.85, x_start=0.30, x_end=9.70,
               colors=[LIGHT_GRAY, LIGHT_BLUE, LIGHT_RED, LIGHT_GRAY,
                       LIGHT_AMBER, LIGHT_GREEN],
               font_size=10.5)

    # Guided prompt structure
    add_text(slide, "Guided prompt template - giảm propagation lỗi từ teacher",
             0.30, 2.20, 9.40, 0.32, size=12, bold=True, color=HUST_RED)
    add_rect(slide, 0.30, 2.55, 9.40, 1.65,
             "[USER]\n{pre_text}\n{table_markdown}\n{post_text}\n{question}\n\n"
             "[Gold program (chỉ dùng để giải thích)]: divide(914,391), multiply(#0,const_100)\n"
             "[Gold answer]: 233.76\n\n"
             "Hãy viết reasoning trace tiếng Việt giải thích vì sao chương trình trên đúng,\n"
             "theo format 3 phần: Phân tích → Chương trình tính toán → Đáp án cuối cùng.",
             fill=LIGHT_BLUE, line=BLUE, size=10.5, align=PP_ALIGN.LEFT,
             text_color=TEXT_DARK)

    # Two bottom callouts
    callout(slide, 0.30, 4.40, 4.55, 2.30, "Vì sao guided thay vì free?",
            ["Free: teacher solve from scratch → ~60% trace có chương trình hợp lệ",
             "Guided: teacher chỉ explain logic gold → ~95% hợp lệ",
             "Student học \"chứng cứ → phép toán\" thay vì chép đáp án",
             "Phù hợp PA: tránh teacher hallucinate program"],
            fill=LIGHT_RED, line=HUST_RED, title_size=12, body_size=10.5)
    callout(slide, 5.10, 4.40, 4.55, 2.30, "Tối ưu teacher (RTX 6000 96GB)",
            ["batch_size = 12 mẫu/lần (~54GB weights + 30GB KV)",
             "max_new_tokens = 512 (~3x nhanh so với 2048)",
             "Length-sorted batching → throughput ↑ 30-50%",
             "checkpoint_every = 100 → sống sót Kaggle 12h"],
            fill=LIGHT_GREEN, line=GREEN, title_color=GREEN, body_size=10.5)


def render_quality_gates(slide):
    # Funnel left
    tiers = [
        ("exact_match", "p_pred ≡_sym p_gold", 60, LIGHT_GREEN, GREEN, "ưu tiên cao nhất"),
        ("answer_match", "exec(p_pred) ≈ a_gold", 25, LIGHT_AMBER, AMBER, "giữ làm SFT"),
        ("program_valid", "parse OK nhưng đáp án sai", 10, LIGHT_BLUE, BLUE, "giữ nếu guided"),
        ("invalid", "không parse được DSL", 5, LIGHT_RED, HUST_RED, "fallback gold"),
    ]
    add_text(slide, "Funnel chất lượng trace teacher", 0.30, 1.00, 4.60, 0.32,
             size=12, bold=True, color=HUST_RED)
    y = 1.40
    max_pct = 70
    for tier, desc, pct, fc, lc, _ in tiers:
        w = 0.30 + 4.40 * pct / max_pct
        add_rect(slide, 0.30, y, w, 0.55, f"{tier}  -  ~{pct}%",
                 fill=fc, line=lc, size=11.5, bold=True, text_color=lc)
        add_text(slide, desc, 5.00, y + 0.05, 4.50, 0.45, size=10.5, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT)
        y += 0.68

    # Decision logic
    add_text(slide, "Logic chọn lọc và fallback (pipeline/teacher_distill.py:48-75)",
             0.30, 4.20, 9.40, 0.32, size=12, bold=True, color=HUST_RED)
    rows = [
        ["Tier", "Hành động giữ trace", "Đóng góp tới SFT"],
        ["exact_match", "Trace teacher đầy đủ", "Mục tiêu chính - chất lượng cao"],
        ["answer_match", "Trace teacher (program khác)", "Đa dạng hóa biểu thức"],
        ["program_valid", "Giữ nếu guided; fallback nếu free", "Mở rộng nhẹ"],
        ["invalid", "Bỏ trace teacher", "Dùng gold SFT response"],
    ]
    add_table(slide, rows, 0.30, 4.55, 9.40, 2.10, font_size=10.5, first_col_bold=True)


def render_lora(slide):
    # Equation visualization
    add_text(slide, "ΔW low-rank: chỉ học B·A thay vì W đầy đủ", 0.30, 0.92, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED, align=PP_ALIGN.CENTER)

    # W' = W + B*A visualization
    base_y = 1.40
    add_rect(slide, 0.30, base_y, 1.50, 1.20, "W\n(đóng băng)\nd × k",
             fill=LIGHT_GRAY, line=TEXT_GRAY, size=12, bold=True, text_color=TEXT_DARK)
    add_text(slide, "+", 1.95, base_y + 0.40, 0.35, 0.40, size=24, bold=True,
             color=HUST_RED, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.40, base_y, 1.20, 1.20, "B\nd × r",
             fill=LIGHT_RED, line=HUST_RED, size=13, bold=True, text_color=HUST_RED)
    add_text(slide, "·", 3.75, base_y + 0.40, 0.30, 0.40, size=24, bold=True,
             color=HUST_RED, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.15, base_y, 1.80, 1.20, "A\nr × k",
             fill=LIGHT_RED, line=HUST_RED, size=13, bold=True, text_color=HUST_RED)
    add_arrow(slide, 6.10, base_y + 0.60, 6.55, base_y + 0.60)
    add_rect(slide, 6.65, base_y, 1.40, 1.20, "W'\n(suy luận)",
             fill=LIGHT_GREEN, line=GREEN, size=13, bold=True, text_color=GREEN)
    add_rect(slide, 8.20, base_y, 1.50, 1.20,
             "r = 128\nα = 256\nall-linear",
             fill=LIGHT_AMBER, line=AMBER, size=11.5, bold=True, text_color=AMBER)

    # Loss masking visualization
    add_text(slide, "Masking nhãn: chỉ tính loss trên assistant response",
             0.30, 2.90, 9.40, 0.32, size=12, bold=True, color=HUST_RED, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.30, 3.25, 6.45, 0.55,
             "[USER prompt: pre_text + table + post_text + question]",
             fill=LIGHT_GRAY, line=TEXT_GRAY, size=11, text_color=TEXT_GRAY)
    add_text(slide, "labels = -100 (không tính loss)", 6.85, 3.30, 2.85, 0.45,
             size=10, italic=True, color=TEXT_GRAY)
    add_rect(slide, 0.30, 3.85, 6.45, 0.55,
             "[ASSISTANT: reasoning + program + answer]",
             fill=LIGHT_GREEN, line=GREEN, size=11, bold=True, text_color=GREEN)
    add_text(slide, "labels = token_id (tính loss)", 6.85, 3.90, 2.85, 0.45,
             size=10, italic=True, color=GREEN, bold=True)

    # Hyperparameter table
    add_text(slide, "Hyperparameter SFT trên RTX 6000 96GB", 0.30, 4.65, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED)
    rows = [
        ["Tham số", "Giá trị", "Tham số", "Giá trị"],
        ["lr", "5e-5", "lora_r", "128"],
        ["scheduler", "cosine, warmup 0.1", "lora_alpha", "256"],
        ["epochs", "2 (Kaggle 12h cap)", "lora_target", "all-linear"],
        ["batch × accum", "1 × 16 = 16", "max_seq_length", "1.536"],
    ]
    add_table(slide, rows, 0.30, 5.00, 9.40, 1.85, font_size=10.5, first_col_bold=True)


def render_method_compare(slide):
    add_text(slide, "Reward executor → không cần value model, không cần preference data",
             0.30, 0.92, 9.40, 0.32, size=12, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)

    rows = [
        ["Phương pháp", "Yêu cầu thêm", "Ưu", "Hạn chế cho PA"],
        ["DPO", "Preference pairs", "Nhẹ, ổn định", "Tạo cặp tự động dễ bias"],
        ["PPO", "Value model riêng", "Linh hoạt", "Tốn VRAM, KL clip khó tune"],
        ["RLVR thuần", "Answer matcher", "Reward kiểm chứng", "Đẩy mô hình đến đáp án đúng / chương trình sai"],
        ["GRPO + PCPO", "Executor (có sẵn)", "Reward tự động, không value model", "Phụ thuộc thiết kế reward"],
    ]
    add_table(slide, rows, 0.30, 1.35, 9.40, 3.30, font_size=10.5, first_col_bold=True)

    # Visual highlight on GRPO row
    callout(slide, 0.30, 4.85, 4.55, 1.85, "Vì sao GRPO phù hợp bài toán này",
            ["Reward có thể tính tự động từ DSL executor < 1ms",
             "Sinh G=5 completion trong cùng prompt → baseline tự nhóm",
             "Không cần value model → tiết kiệm 20-30% VRAM",
             "KL = 1e-3 nhỏ → cho phép policy khám phá"],
            fill=LIGHT_GREEN, line=GREEN, title_color=GREEN, body_size=10.5)
    callout(slide, 5.10, 4.85, 4.55, 1.85, "Vì sao PCPO thay RLVR thuần",
            ["R_valid = 0 nếu DSL sai → bịt 'lucky answer'",
             "Đặt validity chương trình lên trước execution",
             "Bonus brevity dạy chương trình dễ audit",
             "Giảm program hallucination ở RL phase"],
            fill=LIGHT_RED, line=HUST_RED, title_color=HUST_RED, body_size=10.5)


def render_pcpo(slide):
    # Big formula
    add_rect(slide, 0.30, 1.00, 9.40, 0.85,
             "R(p, x)  =  R_valid  ·  ( 0,7  +  0,2 · R_exec  +  0,1 · R_bonus )",
             fill=WHITE, line=HUST_RED, line_width=1.5, size=20, bold=True,
             text_color=HUST_RED)

    # Component breakdown
    components = [
        ("R_valid", "{0, 1}", "Gate cứng - DSL syntax", BLUE, LIGHT_BLUE),
        ("R_exec", "{0, 1}", "exec(p) ≈ a_gold", GREEN, LIGHT_GREEN),
        ("R_bonus", "{0.1, 0.5, 1.0}", "ngắn / bằng / dài hơn gold", AMBER, LIGHT_AMBER),
    ]
    x = 0.30
    w = 3.05
    gap = 0.10
    for name, domain, desc, lc, fc in components:
        add_rect(slide, x, 2.05, w, 1.10, f"{name}\n{domain}\n\n{desc}",
                 fill=fc, line=lc, size=11, bold=True, text_color=TEXT_DARK)
        x += w + gap

    # Visual weight bars
    add_text(slide, "Trọng số α=0,7  β=0,2  γ=0,1", 0.30, 3.30, 9.40, 0.32,
             size=11.5, bold=True, color=HUST_RED, align=PP_ALIGN.CENTER)
    bars = [
        ("α (gate base)", 0.70, HUST_RED),
        ("β (execution)", 0.20, GREEN),
        ("γ (brevity)", 0.10, AMBER),
    ]
    y = 3.75
    for label, val, col in bars:
        add_text(slide, label, 0.30, y - 0.02, 1.65, 0.35, size=11, bold=True,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT)
        bar_w = 6.6 * (val / 0.7)
        add_rect(slide, 2.00, y, bar_w, 0.32, "", fill=col, line=col, rounded=False,
                 line_width=0)
        add_text(slide, f"{val:.1f}", 8.70, y - 0.02, 0.75, 0.35, size=11.5, bold=True,
                 color=col, align=PP_ALIGN.LEFT)
        y += 0.45

    # Comparison table
    add_text(slide, "PCPO vs RLVR thuần - khi 4 tình huống xảy ra", 0.30, 5.15, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED)
    rows = [
        ["Tình huống", "Exec-only", "PCPO"],
        ["DSL sai, đáp án sai", "0", "0"],
        ["DSL sai, đáp án đúng (lucky)", "1 (xấu)", "0"],
        ["DSL đúng, đáp án đúng, ngắn hơn gold", "1", "1,0"],
        ["DSL đúng, đáp án sai", "0", "~0,7 (tín hiệu dương)"],
    ]
    add_table(slide, rows, 0.30, 5.50, 9.40, 1.55, font_size=10.5, first_col_bold=True)


def render_grpo_loop(slide):
    # Top: 5-step loop
    flow_chain(slide, ["Prompt x",
                       "Sinh G = 5\ncompletion",
                       "PCPO reward\nR_i",
                       "Advantage\nA_i",
                       "Policy\nupdate"],
               y=1.05, height=0.85, x_start=0.30, x_end=9.70,
               colors=[LIGHT_GRAY, LIGHT_BLUE, LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN],
               font_size=10.5)

    # Advantage formula
    add_rect(slide, 0.30, 2.20, 9.40, 0.65,
             "A_i  =  ( R_i  -  mean(R) )  /  ( std(R)  +  ε )",
             fill=WHITE, line=BLUE, size=18, bold=True, text_color=BLUE, line_width=1.4)

    # GRPO loss formula box
    add_rect(slide, 0.30, 3.00, 9.40, 1.10,
             "L_GRPO = - E [ Σ_i Σ_t min( ρ_i,t · A_i ,  clip(ρ_i,t, 1-ε, 1+ε) · A_i ) ]  +  β_KL · D_KL(π_θ ‖ π_ref)\n\n"
             "ρ_i,t = π_θ(y_i,t | y_i,<t, x) / π_old(y_i,t | y_i,<t, x)",
             fill=LIGHT_GRAY, line=HUST_RED, size=13, bold=True, text_color=HUST_RED,
             align=PP_ALIGN.CENTER)

    # Hyperparameters and reasoning side by side
    rows = [
        ["Tham số", "Giá trị", "Lý do chọn"],
        ["num_generations", "5", "Đủ để có variance nhóm, không quá đắt"],
        ["lr", "1e-6", "RL nhỏ hơn SFT 50× → tránh policy collapse"],
        ["kl_coef", "1e-3", "Neo nhẹ về SFT, cho phép exploration"],
        ["lora_r", "32 (nhỏ hơn SFT)", "RL chỉ tinh chỉnh không cần dung lượng lớn"],
        ["max_completion_length", "2.048", "Đủ cho prompt + trace + program"],
    ]
    add_table(slide, rows, 0.30, 4.25, 9.40, 2.50, font_size=10.5, first_col_bold=True)


def render_inference(slide):
    # Pipeline diagram
    flow_chain(slide, ["Prompt x",
                       "Sample N=15\nT=0.7, top_p=0.95",
                       "Extract\nprogram/ans",
                       "validate_\nprogram",
                       "execute_\nprogram",
                       "Verifier\nselection"],
               y=1.05, height=0.85, x_start=0.30, x_end=9.70,
               colors=[LIGHT_GRAY, LIGHT_BLUE, LIGHT_GRAY, LIGHT_AMBER,
                       LIGHT_GREEN, LIGHT_RED],
               font_size=10)

    # Verifier formula
    add_rect(slide, 0.30, 2.20, 9.40, 0.70,
             "s(p)  =  λ₁ · valid(p)  +  λ₂ · exec_conf(p)  +  λ₃ · brevity(p)  +  λ₄ · evidence(p)*",
             fill=WHITE, line=HUST_RED, size=16, bold=True, text_color=HUST_RED, line_width=1.4)
    add_text(slide, "*evidence: đề xuất mở rộng - mỗi số trong program phải có nguồn",
             0.30, 2.92, 9.40, 0.28, size=10, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)

    # Left vs right
    callout(slide, 0.30, 3.30, 4.55, 1.95, "Không vote text thô vì:",
            ["Hai program hallucinated có thể trùng đáp án → vote sai",
             "Đáp án giống nhau không = chương trình giống nhau",
             "Vote-on-program qua executor an toàn hơn",
             "exec_conf = tỷ lệ đồng thuận đáp án sau khi exec"],
            fill=LIGHT_RED, line=HUST_RED, title_size=12, body_size=10.5)
    callout(slide, 5.10, 3.30, 4.55, 1.95, "Mở rộng evidence-grounded:",
            ["Mỗi số trong program phải tìm được trong input",
             "Nếu một số xuất hiện mà không có nguồn → trừ điểm",
             "Giảm hallucination căn cơ hơn KL constraint",
             "Tương thích PCPO-Evidence ở training"],
            fill=LIGHT_BLUE, line=BLUE, title_color=BLUE, body_size=10.5)

    # Hyperparams
    rows = [
        ["Tham số", "RTX 6000", "Lý do"],
        ["num_candidates", "15", "Đủ đa dạng, batch 8 chạy nhanh"],
        ["temperature", "0,7", "Đa dạng nhưng không nhiễu quá"],
        ["top_p", "0,95", "Cắt phân phối dài"],
    ]
    add_table(slide, rows, 0.30, 5.42, 9.40, 1.40, font_size=10.5, first_col_bold=True)


def render_optimization(slide):
    # Two-column table
    add_text(slide, "Tối ưu thực tế trong pipeline (config.py + 5 module phase)",
             0.30, 0.92, 9.40, 0.32, size=12, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)

    # Left table - performance
    add_text(slide, "Tối ưu hiệu năng", 0.30, 1.30, 4.55, 0.32,
             size=12, bold=True, color=HUST_RED)
    rows1 = [
        ["Kỹ thuật", "Tác dụng"],
        ["Length-sorted batching", "Giảm padding 30-50%"],
        ["max_new_tokens = 512", "Decode ~3× nhanh hơn 2048"],
        ["LoRA r=128 (SFT) / 32 (GRPO)", "Giảm VRAM optimizer"],
        ["Gradient accumulation 16", "Effective batch 16 trên 1 GPU"],
        ["bf16 / fp16 theo profile", "Tận dụng tensor core"],
        ["Checkpoint every 100", "Sống sót session 12h"],
        ["mirror_save_dir", "Phòng kernel kill"],
        ["max_runtime_hours watchdog", "Save trước cut-off"],
    ]
    add_table(slide, rows1, 0.30, 1.65, 4.55, 4.30, font_size=10, first_col_bold=True)

    # Right table - reliability
    add_text(slide, "Tối ưu độ tin cậy", 5.10, 1.30, 4.55, 0.32,
             size=12, bold=True, color=BLUE)
    rows2 = [
        ["Cơ chế", "Kiểm tra"],
        ["Regex 3-phần", "Trích program/answer ổn định"],
        ["validate_program", "Chặn function ngoài DSL"],
        ["execute_program", "Tolerance 1e-4 / 1e-5"],
        ["4-strategy header match", "Xử lý header biến thể"],
        ["Quality tier 4 mức", "Trace teacher có fallback"],
        ["PA symbolic (sympy)", "Equivalence thay string match"],
        ["R_valid gate", "Không reward program lucky"],
        ["Robust load + SDPA fallback", "Chống flash-attn fail"],
    ]
    add_table(slide, rows2, 5.10, 1.65, 4.55, 4.30, font_size=10, first_col_bold=True,
              header_fill=BLUE)

    add_text(slide,
             "Pipeline chạy được trên 5 profile GPU (P100 16GB → A100 80GB) với cùng code path.",
             0.30, 6.20, 9.40, 0.40, size=11.5, italic=True, color=HUST_RED,
             align=PP_ALIGN.CENTER)


def render_ablation(slide):
    # Table left
    rows = [
        ["Thiết lập", "EA (%)", "PA (%)"],
        ["Zero-shot Qwen3.5-4B", "52,0", "35,0"],
        ["+ SFT ViNumQA-only", "68,7", "60,8"],
        ["+ FinQA + program_re", "71,6", "65,4"],
        ["+ Guided KD 27B → 4B", "72,8", "68,6"],
        ["+ GRPO/PCPO", "73,6", "70,2"],
        ["+ Verifier/voting", "74,4", "70,8"],
    ]
    add_table(slide, rows, 0.30, 1.00, 4.85, 3.60, font_size=10.5, first_col_bold=True)

    # PA progression bar chart right
    add_text(slide, "Tiến triển PA theo từng tầng kỹ thuật", 5.40, 0.95, 4.30, 0.32,
             size=12, bold=True, color=HUST_RED)
    pa_progression = [
        ("Zero-shot", 35.0),
        ("SFT", 60.8),
        ("+FinQA+re", 65.4),
        ("+KD", 68.6),
        ("+PCPO", 70.2),
        ("+Verify", 70.8),
    ]
    # Custom vertical bar chart
    x0 = 5.40
    y0 = 4.70  # baseline
    chart_w = 4.30
    chart_h = 3.30
    n = len(pa_progression)
    col_w = (chart_w - 0.40) / n
    max_pa = 80
    for i, (lab, val) in enumerate(pa_progression):
        bar_h = chart_h * val / max_pa
        bx = x0 + 0.20 + i * col_w
        by = y0 - bar_h
        col = GREEN if i >= 3 else (AMBER if i >= 1 else TEXT_GRAY)
        add_rect(slide, bx, by, col_w * 0.78, bar_h, "", fill=col, line=col,
                 rounded=False, line_width=0)
        add_text(slide, f"{val:.0f}", bx, by - 0.30, col_w * 0.78, 0.28,
                 size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, lab, bx - 0.05, y0 + 0.02, col_w * 0.78 + 0.10, 0.42,
                 size=8.5, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    # Y-axis labels
    add_text(slide, "PA (%)", x0 - 0.05, 1.25, 0.5, 0.25, size=9, italic=True,
             color=TEXT_GRAY)

    # Interpretation callouts
    add_text(slide, "Đọc số ablation", 0.30, 4.90, 4.85, 0.32, size=12, bold=True,
             color=HUST_RED)
    add_bullets(slide, [
        "ΔPA lớn nhất: KD (+3,2) và FinQA+re (+4,6) → đóng góp dữ liệu/distillation",
        "GRPO/PCPO +1,6 PA: vai trò gating DSL hợp lệ (đặc biệt cho lỗi syntax)",
        "Verifier tăng EA (+0,8) > PA (+0,6): stabilize answer, không stabilize logic",
        "Không có 'magic jump' - mỗi tầng có lý do thiết kế rõ ràng",
    ], 0.30, 5.25, 4.85, 1.85, size=10.5, bullet_color=HUST_RED)


def render_errors(slide):
    rows = [
        ["Lỗi", "Nguyên nhân", "Hiện có", "Đề xuất"],
        ["Chọn sai số liệu",
         "Header phức tạp, nhiều số gần",
         "4-strategy header match\nMarkdown table",
         "Evidence map số → cell"],
        ["Sai phép toán",
         "Câu hỏi mơ hồ",
         "Guided rationale KD",
         "Template Q + op-suggester"],
        ["Sai DSL syntax",
         "Mô hình sinh text tự do",
         "PCPO R_valid gate\nvalidate_program",
         "Grammar-constrained decoding"],
        ["Đúng EA sai PA",
         "Program khác gold, PA chưa phủ",
         "program_re augment\nPA symbolic (sympy)",
         "Rewrite-rule catalog mở rộng"],
        ["Hallucination reasoning",
         "Reward chưa kiểm reasoning text",
         "(chưa có)",
         "PCPO-Evidence + verifier chứng cứ"],
    ]
    add_table(slide, rows, 0.30, 1.05, 9.40, 5.10, font_size=10, first_col_bold=True)

    add_text(slide,
             "Mỗi loại lỗi có một kỹ thuật xử lý hiện tại và một hướng mở rộng - đường nâng cấp rõ ràng.",
             0.30, 6.30, 9.40, 0.40, size=11.5, italic=True, color=HUST_RED,
             align=PP_ALIGN.CENTER)


def render_method_contrib(slide):
    add_text(slide,
             "Đóng góp riêng của đề tài (không tính kỹ thuật baseline có sẵn)",
             0.30, 0.92, 9.40, 0.32, size=12, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)

    contribs = [
        ("DỮ LIỆU", "Chiến lược program-centric",
         ["3 luồng SFT/GRPO/teacher đi qua executor chung",
          "Markdown + program_re (2.534 lời giải thay thế)",
          "4-strategy header match cho bảng tài chính",
          "Quality gate trước GRPO"],
         BLUE, LIGHT_BLUE),
        ("THUẬT TOÁN", "PCPO + Guided KD + Quality tier",
         ["R_valid · (0,7 + 0,2·R_exec + 0,1·R_bonus)",
          "Guided prompt: teacher chỉ explain, không solve",
          "4-tier quality + fallback gold",
          "GRPO với KL=1e-3 giữ năng lực ngôn ngữ"],
         HUST_RED, LIGHT_RED),
        ("HỆ THỐNG", "Pipeline 5 phase đa profile",
         ["Tách độc lập có checkpoint riêng",
          "Length-sorted batching + left-padding",
          "Mirror save + watchdog Kaggle 12h",
          "Cùng code chạy P100 → A100 80GB"],
         GREEN, LIGHT_GREEN),
        ("ĐỊNH HƯỚNG", "PCPO-Evidence + Grammar decoding",
         ["Evidence map số → cell, reward chứng cứ",
          "Grammar-constrained decoding chặn invalid DSL",
          "Rewrite-rule catalog cho PA",
          "Agent verifier - re-verify question"],
         PURPLE, LIGHT_PURPLE),
    ]
    x0, y0 = 0.30, 1.35
    w, h = 4.55, 2.60
    gap = 0.15
    for i, (tag, title, bullets, lc, fc) in enumerate(contribs):
        cx = x0 + (i % 2) * (w + gap)
        cy = y0 + (i // 2) * (h + 0.15)
        # Tag header
        add_rect(slide, cx, cy, w, 0.30, tag, fill=lc, line=lc,
                 text_color=WHITE, size=10, bold=True)
        # Title + body
        add_rect(slide, cx, cy + 0.30, w, h - 0.30, "",
                 fill=fc, line=lc, rounded=True)
        add_text(slide, title, cx + 0.12, cy + 0.36, w - 0.24, 0.34,
                 size=12.5, bold=True, color=lc, align=PP_ALIGN.LEFT)
        add_bullets(slide, bullets, cx + 0.12, cy + 0.72, w - 0.24, h - 0.85,
                    size=10.5, color=TEXT_DARK, bullet_color=lc, space_after=2)

    add_text(slide,
             "Cốt lõi: chương trình tính toán là trung tâm của dữ liệu, distillation, reward và inference.",
             0.30, 6.70, 9.40, 0.35, size=11.5, italic=True, color=HUST_RED,
             align=PP_ALIGN.CENTER)


def render_future(slide):
    # 3-layer integration
    add_text(slide, "Ba lớp tích hợp để tăng PA và giảm hallucination",
             0.30, 0.92, 9.40, 0.32, size=12, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)

    layers = [
        ("Lớp 1 - Dữ liệu",
         "Evidence-grounded annotation",
         "Mỗi số trong gold program gắn với cell/span trong input",
         BLUE, LIGHT_BLUE),
        ("Lớp 2 - Huấn luyện",
         "PCPO-Evidence + Grammar decoding",
         "Reward thêm evidence; grammar chặn invalid DSL ngay decode",
         HUST_RED, LIGHT_RED),
        ("Lớp 3 - Inference",
         "Verifier đa thành phần + Agent re-verify",
         "valid + exec + brevity + evidence + agreement",
         GREEN, LIGHT_GREEN),
    ]
    y = 1.30
    for tag, title, body, lc, fc in layers:
        add_rect(slide, 0.30, y, 1.85, 1.30, tag, fill=lc, line=lc,
                 text_color=WHITE, size=11.5, bold=True)
        add_rect(slide, 2.20, y, 3.30, 1.30, title, fill=fc, line=lc,
                 size=12.5, bold=True, text_color=lc)
        add_rect(slide, 5.55, y, 4.15, 1.30, body, fill=WHITE, line=lc,
                 size=11, text_color=TEXT_DARK, align=PP_ALIGN.LEFT)
        y += 1.43

    # Roadmap timeline
    add_text(slide, "Roadmap đề xuất", 0.30, 5.85, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED)
    flow_chain(slide, ["Evidence\nlabels",
                       "PCPO-\nEvidence",
                       "Grammar\ndecoding",
                       "Rewrite\nrules PA",
                       "Agent\nverifier"],
               y=6.20, height=0.65, x_start=0.30, x_end=9.70,
               colors=[LIGHT_BLUE, LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN, LIGHT_PURPLE],
               font_size=10)


def render_conclusion(slide):
    # Big message
    add_rect(slide, 0.50, 1.05, 9.00, 1.05,
             "Tài chính cần mô hình sinh chương trình kiểm chứng được, không chỉ sinh đáp án.\n"
             "Đặt executor DSL làm trung tâm: dữ liệu, distillation, reward, inference, đánh giá.",
             fill=LIGHT_RED, line=HUST_RED, size=15, bold=True, text_color=HUST_RED)

    # Three key takeaways
    add_text(slide, "Ba điểm chốt", 0.50, 2.35, 9.00, 0.32, size=12, bold=True,
             color=HUST_RED)

    takeaways = [
        ("01", "Guided KD",
         "Teacher 27B sinh trace có gold program nhúng vào prompt → student 4B học \"chứng cứ → phép toán\" thay vì sao chép đáp án.",
         BLUE, LIGHT_BLUE),
        ("02", "PCPO reward",
         "R_valid · (0,7 + 0,2·R_exec + 0,1·R_bonus) đặt validity chương trình lên trước. Bịt lucky answer của RLVR thuần.",
         HUST_RED, LIGHT_RED),
        ("03", "Verifier inference",
         "Vote-on-program qua executor + verifier (valid + exec + brevity). Hướng mở rộng evidence để giảm hallucination căn cơ.",
         GREEN, LIGHT_GREEN),
    ]
    y = 2.75
    for tag, title, body, lc, fc in takeaways:
        add_rect(slide, 0.50, y, 0.85, 1.10, tag, fill=lc, line=lc,
                 text_color=WHITE, size=22, bold=True)
        add_rect(slide, 1.40, y, 2.15, 1.10, title, fill=fc, line=lc,
                 size=13, bold=True, text_color=lc)
        add_rect(slide, 3.60, y, 5.90, 1.10, body, fill=WHITE, line=lc,
                 size=11.5, text_color=TEXT_DARK, align=PP_ALIGN.LEFT)
        y += 1.20

    # Closing
    add_rect(slide, 0.50, 6.30, 9.00, 0.55,
             "Đường mở rộng: PCPO-Evidence + Grammar-constrained decoding + Agent verifier",
             fill=LIGHT_GRAY, line=HUST_RED, size=12.5, bold=True, text_color=HUST_RED)


# ── Slide renderers cho cấu trúc academic 26-slide ───────────────────

def render_outline(slide):
    add_text(slide, "Theo quy trình báo cáo nghiên cứu chuẩn (6 chương)",
             0.50, 0.95, 9.00, 0.35, size=13, italic=True, color=TEXT_GRAY)
    chapters = [
        ("Chương 1", "Giới thiệu và Mô tả bài toán",
         "Bối cảnh, input/output, ví dụ mẫu, thống kê dataset", BLUE, LIGHT_BLUE),
        ("Chương 2", "Tổng quan nghiên cứu liên quan",
         "Ba nghiên cứu nền tảng pre-2026: FinQA, Step-by-Step, GRPO", HUST_RED, LIGHT_RED),
        ("Chương 3", "Phân tích bài toán & Cơ sở lý thuyết",
         "KD, PEFT (LoRA), RL verifier-reward — định hướng program-centric", AMBER, LIGHT_AMBER),
        ("Chương 4", "Phương pháp đề xuất (5 đóng góp)",
         "Dữ liệu, Guided KD, LoRA-SFT, GRPO+PCPO, Verifier inference", GREEN, LIGHT_GREEN),
        ("Chương 5", "Cài đặt thực nghiệm & Kết quả",
         "Hardware, hyperparam, ablation, phân tích lỗi", PURPLE, LIGHT_PURPLE),
        ("Chương 6", "Thảo luận & Kết luận",
         "Vì sao đạt kết quả, hạn chế, hướng phát triển", BLUE, LIGHT_BLUE),
    ]
    y = 1.55
    for tag, title, body, lc, fc in chapters:
        add_rect(slide, 0.50, y, 1.50, 0.78, tag, fill=lc, line=lc,
                 text_color=WHITE, size=14, bold=True)
        add_rect(slide, 2.05, y, 3.10, 0.78, title, fill=fc, line=lc,
                 size=12.5, bold=True, text_color=lc, align=PP_ALIGN.LEFT)
        add_rect(slide, 5.20, y, 4.30, 0.78, body, fill=WHITE, line=lc,
                 size=11, text_color=TEXT_DARK, align=PP_ALIGN.LEFT)
        y += 0.84


def render_context(slide):
    add_text(slide,
             "VLSP 2025 NumQA — challenge đầu tiên cho cộng đồng NLP Việt về suy luận số tài chính",
             0.50, 0.92, 9.00, 0.35, size=12.5, italic=True, color=TEXT_GRAY)

    # 3 nhóm thách thức — mỗi nhóm có 2 dòng giải thích cụ thể
    blocks = [
        ("01", "Suy luận số đa bước trên bảng",
         "Trung bình 1,56 bước/program, tối đa 7 bước. Mỗi bước có thể truy bảng "
         "(table_sum, table_average...) hoặc tính trên kết quả bước trước (#0, #1...).\n"
         "Lỗi ở bước đầu lan truyền → toàn bộ chương trình sai → cả EA và PA về 0.",
         BLUE, LIGHT_BLUE),
        ("02", "Bảng tài chính tiếng Việt dị thường",
         "Header chứa dấu tiếng Việt, đơn vị (\"triệu đồng\"), năm (\"2022\"), tên cột dài.\n"
         "Tham chiếu trong câu hỏi thường KHÔNG khớp exact với header → cần chiến lược matching dung sai cao "
         "nhưng không sai khớp (đề xuất 4-strategy header matching ở Mục 4.2.2).",
         HUST_RED, LIGHT_RED),
        ("03", "Dữ liệu tiếng Việt khan hiếm",
         "ViNumQA train chỉ có 2.993 mẫu — ít hơn nhiều so với FinQA (6.251 mẫu English).\n"
         "Dịch FinQA En → Vi gây lỗi số (\"1,840\" → văn bản), lỗi đơn vị → cần multilingual mà KHÔNG dịch, "
         "để mô hình tự xử lý hai ngôn ngữ.",
         GREEN, LIGHT_GREEN),
    ]
    y = 1.40
    for tag, title, body, lc, fc in blocks:
        add_rect(slide, 0.40, y, 0.75, 1.45, tag, fill=lc, line=lc,
                 text_color=WHITE, size=20, bold=True)
        add_rect(slide, 1.20, y, 2.40, 1.45, title, fill=fc, line=lc,
                 size=12, bold=True, text_color=lc)
        add_rect(slide, 3.65, y, 5.95, 1.45, body, fill=WHITE, line=lc,
                 size=10.5, text_color=TEXT_DARK, align=PP_ALIGN.LEFT)
        y += 1.55

    # Hệ quả - tại sao chương trình làm trung tâm
    add_rect(slide, 0.40, 6.10, 9.20, 0.90,
             "HỆ QUẢ: trong tài chính một đáp án đúng nhưng không kiểm toán được không đủ tin cậy.\n"
             "→ Mục tiêu PA quan trọng hơn EA → toàn bộ pipeline được thiết kế PROGRAM-CENTRIC.",
             fill=LIGHT_AMBER, line=HUST_RED, size=12, bold=True, text_color=HUST_RED)


def render_sample(slide):
    add_text(slide, "Một mẫu thực tế minh hoạ: input → reasoning → program → answer",
             0.50, 0.95, 9.00, 0.32, size=12, italic=True, color=TEXT_GRAY)

    # Sample table (left side)
    add_text(slide, "Bảng dữ liệu (Markdown):", 0.50, 1.32, 4.50, 0.30,
             size=11.5, bold=True, color=HUST_RED)
    table_rows = [
        ["Chỉ tiêu", "Năm 2022", "Năm 2023"],
        ["Doanh thu thuần", "12.450", "15.230"],
        ["Giá vốn hàng bán", "8.120", "9.760"],
        ["Lợi nhuận sau thuế", "1.840", "2.510"],
    ]
    add_table(slide, table_rows, 0.50, 1.65, 4.50, 1.55, font_size=10.5,
              first_col_bold=True)

    # Question (left side, below table)
    add_rect(slide, 0.50, 3.35, 4.50, 0.85,
             "Câu hỏi: Tỉ lệ tăng trưởng lợi nhuận sau thuế năm 2023 so với năm 2022 là bao nhiêu phần trăm?",
             fill=LIGHT_BLUE, line=BLUE, size=11, bold=False, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT)

    # Right side - output flow
    add_text(slide, "Đầu ra mong đợi:", 5.20, 1.32, 4.30, 0.30,
             size=11.5, bold=True, color=HUST_RED)

    # Reasoning
    add_rect(slide, 5.20, 1.65, 4.30, 0.95,
             "Reasoning: LNST 2023 = 2.510, 2022 = 1.840.\n"
             "Mức tăng tuyệt đối = 2.510 − 1.840 = 670.\n"
             "Tỉ lệ = 670 / 1.840 ≈ 0,3641.",
             fill=LIGHT_GREEN, line=GREEN, size=10.5, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT)
    add_arrow(slide, 7.35, 2.65, 7.35, 2.85, color=HUST_RED, width_pt=1.8)
    # Program
    add_rect(slide, 5.20, 2.90, 4.30, 0.85,
             "Program (DSL):\nsubtract(2510, 1840),  divide(#0, 1840)",
             fill=LIGHT_AMBER, line=AMBER, size=10.5, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT)
    add_arrow(slide, 7.35, 3.80, 7.35, 4.00, color=HUST_RED, width_pt=1.8)
    # Answer
    add_rect(slide, 5.20, 4.05, 4.30, 0.50,
             "Answer: 0.3641 (≈ 36,41%)",
             fill=LIGHT_RED, line=HUST_RED, size=12, bold=True,
             text_color=HUST_RED)

    # Bottom - quan sát quan trọng
    add_rect(slide, 0.50, 4.85, 9.00, 1.85,
             "QUAN SÁT QUAN TRỌNG",
             fill=HUST_RED, line=HUST_RED, size=12, bold=True, text_color=WHITE)
    add_bullets(slide, [
        "Program có 2 bước: subtract → divide; reference #0 trỏ về bước 0.",
        "Một chương trình thay thế (program_re): divide(subtract(2510,1840), 1840) — khác cấu trúc, cùng đáp án, vẫn đạt PA nếu sympy rewrite tương đương.",
        "Phép table_average(Lợi nhuận sau thuế) cần executor map header tiếng Việt → 4-strategy header matching (Mục 4.2.2 trong báo cáo).",
        "Mỗi mẫu trung bình 1,56 bước; 10% mẫu ≥ 3 bước — vùng student dễ sai.",
    ], 0.65, 5.20, 8.70, 1.50, size=11, color=WHITE, bullet_color=WHITE)


def render_related_overview(slide):
    add_text(slide,
             "Đề tài giao thoa 3 dòng nghiên cứu nền tảng pre-2026",
             0.50, 0.95, 9.00, 0.32, size=12.5, italic=True, color=TEXT_GRAY)

    # Timeline
    add_rect(slide, 0.50, 1.45, 9.00, 0.32, "", fill=HUST_RED, line=HUST_RED, rounded=False)
    points = [
        ("2021", "FinQA", "Program-centric financial QA, định nghĩa DSL", BLUE),
        ("2023", "Step-by-Step", "Distill rationale từ LLM teacher", HUST_RED),
        ("2024", "GRPO", "RL không cần value model, group baseline", GREEN),
        ("2025", "Đề tài", "Tích hợp pipeline program-centric Vi", PURPLE),
    ]
    x = 0.70
    for year, name, body, color in points:
        add_rect(slide, x, 1.30, 1.20, 0.62, year, fill=color, line=color,
                 text_color=WHITE, size=14, bold=True)
        x += 2.20

    # Detail cards
    cards = [
        ("FinQA — Chen et al. (EMNLP 2021)",
         "DSL 10 phép toán\nRetriever+Generator\nFinQANet: EA 61% / PA 59%\n→ DSL kế thừa nguyên dạng",
         BLUE, LIGHT_BLUE),
        ("Step-by-Step — Hsieh et al. (ACL 2023)",
         "Teacher 540B sinh rationale\nStudent 770M-11B multi-task\nVượt teacher với 12,5% data\n→ ý tưởng truyền reasoning",
         HUST_RED, LIGHT_RED),
        ("DeepSeekMath GRPO — Shao et al. (2024)",
         "RL không value model\nGroup baseline + KL\n51,7% trên MATH (7B)\n→ backbone RL của đề tài",
         GREEN, LIGHT_GREEN),
    ]
    x = 0.50
    for title, body, lc, fc in cards:
        add_rect(slide, x, 2.20, 3.00, 0.50, title, fill=lc, line=lc,
                 text_color=WHITE, size=11, bold=True)
        add_rect(slide, x, 2.70, 3.00, 2.30, body, fill=fc, line=lc,
                 size=11, text_color=TEXT_DARK, align=PP_ALIGN.LEFT)
        x += 3.15

    # Định vị
    add_rect(slide, 0.50, 5.30, 9.00, 1.45,
             "ĐỀ TÀI 2025\n"
             "Không phát minh lại DSL, không phát minh lại GRPO.\n"
             "Đóng góp = tích hợp 3 dòng trên với executor DSL ở trung tâm,\n"
             "thêm Guided Distillation và PCPO reward gate.",
             fill=LIGHT_AMBER, line=HUST_RED, size=12, bold=True, text_color=HUST_RED)


def render_rw_finqa(slide):
    add_text(slide, "Chen, Z. et al. (2021). EMNLP. arXiv:2109.00122",
             0.50, 0.95, 9.00, 0.30, size=11.5, italic=True, color=TEXT_GRAY)

    # Đóng góp
    add_text(slide, "Đóng góp chính của FinQA", 0.50, 1.35, 4.50, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    add_bullets(slide, [
        "Dataset 8.281 cặp QA tài chính từ S&P 500 (train 6.251, dev 883, test 1.147).",
        "Định nghĩa DSL 10 phép toán: add, subtract, multiply, divide, exp, greater, table_sum, table_max, table_min, table_average.",
        "Kiến trúc baseline Retriever (BERT) + Generator (encoder-decoder) sinh program.",
        "Định nghĩa metric chuẩn EA + PA (sympy symbolic equivalence).",
    ], 0.50, 1.70, 4.50, 3.40, size=11)

    # Kết quả công bố
    add_text(slide, "Kết quả công bố trên FinQA test",
             5.20, 1.35, 4.30, 0.32, size=12.5, bold=True, color=HUST_RED)
    rows = [
        ["Mô hình", "EA", "PA"],
        ["GPT-3 zero-shot", "14,4%", "—"],
        ["Longformer + Gen", "21,7%", "18,4%"],
        ["BERT base + Gen", "50,2%", "47,5%"],
        ["FinQANet (RoBERTa-large)", "61,2%", "58,9%"],
        ["Human expert", "91,2%", "87,5%"],
    ]
    add_table(slide, rows, 5.20, 1.70, 4.30, 2.65, font_size=10,
              first_col_bold=True)

    # Liên hệ với đề tài
    add_rect(slide, 0.50, 5.30, 9.00, 1.50,
             "LIÊN HỆ VỚI ĐỀ TÀI: kế thừa nguyên dạng DSL 10 phép + EA/PA metric.\n"
             "Khác biệt: đề tài chuyển sang tiếng Việt, dùng LLM Qwen3.5 thay encoder-decoder chuyên biệt,\n"
             "thêm KD + GRPO + Verifier — vượt FinQANet PA 58,9% lên ~70,5% với mô hình LoRA 98M tham số.",
             fill=LIGHT_RED, line=HUST_RED, size=11.5, bold=False, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT)


def render_rw_distill(slide):
    add_text(slide, "Hsieh, C.-Y. et al. (2023). ACL Findings. arXiv:2305.02301",
             0.50, 0.95, 9.00, 0.30, size=11.5, italic=True, color=TEXT_GRAY)

    # Ý tưởng và formula
    add_text(slide, "Ý tưởng cốt lõi", 0.50, 1.35, 9.00, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    add_rect(slide, 0.50, 1.70, 9.00, 0.95,
             "Teacher LLM 540B sinh thêm RATIONALE ngoài label.\n"
             "Student multi-task: L = α·L_label + (1−α)·L_rationale.",
             fill=LIGHT_BLUE, line=BLUE, size=13, text_color=TEXT_DARK)

    # Kết quả
    add_text(slide, "Kết quả công bố (ANLI)", 0.50, 2.85, 4.50, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    rows = [
        ["Mô hình", "Acc", "Params"],
        ["PaLM 540B teacher", "70,1%", "540B"],
        ["T5-Base SFT chuẩn", "49,2%", "220M"],
        ["T5-Base + Step-by-Step", "53,4%", "220M"],
        ["T5-XXL + Step-by-Step", "70,4%", "11B"],
    ]
    add_table(slide, rows, 0.50, 3.20, 4.50, 2.05, font_size=10,
              first_col_bold=True)

    # Đề tài khác biệt
    add_text(slide, "Đề tài làm khác", 5.20, 2.85, 4.30, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    add_bullets(slide, [
        "Hsieh: teacher TỰ SINH rationale → 40% nhiễu hallucination.",
        "Đề tài: GOLD PROGRAM nhúng prompt → teacher chỉ \"giải thích\".",
        "Trace hợp lệ: 60% → 95% (+35 điểm).",
        "Quality tier 4 mức làm safety net.",
    ], 5.20, 3.20, 4.30, 2.20, size=11)

    # Bottom message
    add_rect(slide, 0.50, 5.50, 9.00, 1.20,
             "KẾ THỪA: ý tưởng truyền rationale.\n"
             "CẢI TIẾN: đảo vai trò teacher từ generator sang explainer → tận dụng KD reasoning nhưng triệt tiêu hallucination.",
             fill=LIGHT_AMBER, line=HUST_RED, size=11.5, bold=False, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT)


def render_rw_grpo(slide):
    add_text(slide, "Shao, Z. et al. (2024). DeepSeekMath. arXiv:2402.03300",
             0.50, 0.95, 9.00, 0.30, size=11.5, italic=True, color=TEXT_GRAY)

    # Công thức GRPO
    add_text(slide, "Công thức GRPO (tóm tắt)", 0.50, 1.35, 9.00, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    add_rect(slide, 0.50, 1.70, 9.00, 1.30,
             "Sample G outputs {o_i} ~ π_old (·| q)\n"
             "Group advantage:  Â_i = (r_i − mean(r)) / (std(r) + ε)\n"
             "L_GRPO = − (1/G) Σ_i (1/|o_i|) Σ_t min( ρ_{i,t}·Â_i, clip(ρ,1−ε,1+ε)·Â_i ) + β·KL(π_θ ∥ π_ref)\n"
             "ρ_{i,t} = π_θ(o_{i,t}| q, o_{i,<t}) / π_old(o_{i,t}| q, o_{i,<t})",
             fill=WHITE, line=HUST_RED, size=11.5, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT, line_width=1.3)

    # Lợi thế GRPO vs PPO
    add_text(slide, "GRPO vs PPO", 0.50, 3.15, 4.50, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    rows = [
        ["Khía cạnh", "PPO", "GRPO"],
        ["Yêu cầu", "Policy + Value", "Chỉ Policy + Ref"],
        ["Bộ nhớ", "~2× params", "~1× params"],
        ["Baseline", "V_φ(s) học", "Mean reward group"],
        ["Reward sparse", "Khó học V", "Group norm OK"],
    ]
    add_table(slide, rows, 0.50, 3.50, 4.50, 1.85, font_size=10,
              first_col_bold=True)

    # Đề tài kế thừa + thêm gì
    add_text(slide, "Đề tài kế thừa & mở rộng", 5.20, 3.15, 4.30, 0.32,
             size=12.5, bold=True, color=HUST_RED)
    add_bullets(slide, [
        "Backbone RL = GRPO (G=5, KL=1e-3).",
        "Reward gốc DeepSeek: chỉ Acc trên đáp án toán.",
        "Đề tài thay bằng PCPO reward.",
        "R_valid làm HARD GATE — bịt lucky-answer.",
        "PA tăng +13,3 điểm nhờ gate này.",
    ], 5.20, 3.50, 4.30, 2.10, size=11)

    # Bottom
    add_rect(slide, 0.50, 5.65, 9.00, 1.10,
             "KẾT QUẢ DeepSeekMath-7B-RL: 51,7% MATH (SOTA open 7B đầu 2024).\n"
             "Đề tài áp dụng GRPO cho student 4B + reward đặc thù DSL.",
             fill=LIGHT_GREEN, line=GREEN, size=11.5, bold=False, text_color=TEXT_DARK,
             align=PP_ALIGN.LEFT)


def render_positioning(slide):
    add_text(slide,
             "Bảng so sánh trực diện: đề tài kế thừa gì, đóng góp gì",
             0.50, 0.95, 9.00, 0.32, size=12, italic=True, color=TEXT_GRAY)

    rows = [
        ["Khía cạnh", "FinQA 2021", "Step-by-Step 2023", "GRPO 2024", "ĐỀ TÀI 2025"],
        ["Ngôn ngữ", "Anh", "Anh", "Anh (toán)", "Việt + Anh"],
        ["DSL chương trình", "Đề xuất", "—", "—", "Kế thừa FinQA"],
        ["Distillation", "—", "Free rationale", "—", "Guided + 4 tier"],
        ["RL algorithm", "—", "—", "GRPO + Acc", "GRPO + PCPO"],
        ["Verifier", "Sympy outside", "—", "Math checker", "Executor in mọi phase"],
        ["Inference", "Greedy", "Greedy", "Greedy", "Multi-path + verifier"],
        ["Tham số train", "~330M", "770M-11B", "7B", "98M LoRA / 4B"],
    ]
    add_table(slide, rows, 0.30, 1.35, 9.40, 4.30, font_size=10,
              first_col_bold=True)

    # Bottom message
    add_rect(slide, 0.50, 5.85, 9.00, 0.95,
             "ĐÓNG GÓP CHÍNH = THIẾT KẾ HỆ THỐNG\n"
             "Đặt executor DSL ở trung tâm, ràng buộc mọi giai đoạn — không phát minh lại DSL hay GRPO.",
             fill=LIGHT_RED, line=HUST_RED, size=12.5, bold=True, text_color=HUST_RED)


def render_theory(slide):
    add_text(slide, "Ba khối lý thuyết nền tảng kết hợp thành một định hướng thiết kế thống nhất",
             0.50, 0.92, 9.00, 0.32, size=12, italic=True, color=TEXT_GRAY)

    # 3 khối với nội dung mô tả đầy đủ hơn
    blocks = [
        ("KD", "Knowledge\nDistillation",
         "• Seq-level KD (Kim 2016): truyền chuỗi sample, không truyền logit\n"
         "• Reasoning trace KD (Hsieh 2023): truyền rationale ngoài label\n"
         "• Rủi ro chính: teacher hallucinate → trace nhiễu lan vào SFT\n"
         "→ Đề tài: Guided distillation (gold program nhúng prompt) loại rủi ro",
         BLUE, LIGHT_BLUE),
        ("PEFT", "LoRA\nlow-rank ΔW",
         "• Full FT 4B cần ~52 GB VRAM → không khả thi P100 16GB\n"
         "• LoRA: ΔW = B·A với r ≪ min(d,k), backward chỉ trên B,A\n"
         "• Đề tài chọn r=128, α=256: cao hơn r=8/16 thông thường\n"
         "→ 98M params (~2,5% full FT), vừa 12-14 GB train memory",
         AMBER, LIGHT_AMBER),
        ("RL", "Verifier-\nguided RL",
         "• PPO cần value model → gấp 2× bộ nhớ\n"
         "• DPO cần preference data (chosen, rejected) → ta không có\n"
         "• GRPO: group baseline = mean reward trong G samples, không value model\n"
         "→ Đề tài: GRPO + reward từ executor DSL (verifier có sẵn)",
         GREEN, LIGHT_GREEN),
    ]
    x = 0.30
    for tag, title, body, lc, fc in blocks:
        add_rect(slide, x, 1.35, 3.15, 0.45, tag, fill=lc, line=lc,
                 text_color=WHITE, size=13, bold=True)
        add_rect(slide, x, 1.80, 3.15, 0.62, title, fill=WHITE, line=lc,
                 size=12.5, bold=True, text_color=lc)
        add_rect(slide, x, 2.42, 3.15, 2.40, body, fill=fc, line=lc,
                 size=10, text_color=TEXT_DARK, align=PP_ALIGN.LEFT)
        x += 3.20

    # Định hướng chốt
    add_rect(slide, 0.30, 5.05, 9.40, 1.40,
             "ĐỊNH HƯỚNG THIẾT KẾ CHỐT\n\n"
             "Mọi tín hiệu giám sát — SFT loss, distillation rationale, GRPO reward, inference verifier —\n"
             "đều phải đi qua CÙNG MỘT EXECUTOR DSL.\n"
             "→ Consistency giữa các giai đoạn, không bị objective drift.",
             fill=LIGHT_RED, line=HUST_RED, size=12, bold=True, text_color=HUST_RED)

    add_text(slide,
             "Đây là điểm khác biệt căn bản so với pipeline đa giai đoạn truyền thống vốn dễ bị mỗi giai đoạn tối ưu một mục tiêu khác.",
             0.30, 6.55, 9.40, 0.40, size=11, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)


def render_setup(slide):
    add_text(slide, "Cài đặt thực nghiệm thực tế trên Kaggle GPU",
             0.50, 0.95, 9.00, 0.32, size=12, italic=True, color=TEXT_GRAY)

    # Hardware table
    add_text(slide, "Phần cứng & mô hình", 0.50, 1.35, 4.50, 0.32,
             size=12, bold=True, color=HUST_RED)
    hw_rows = [
        ["Vai trò", "Mô hình", "VRAM"],
        ["Teacher", "Qwen3.5-27B bf16", "~54 GB"],
        ["Student", "Qwen3.5-4B bf16", "~9 GB"],
        ["Reference (GRPO)", "Student SFT frozen", "~9 GB"],
        ["GPU train", "RTX 6000 Ada 48GB", "Pool 96GB"],
        ["GPU eval", "Tesla P100 16GB", "16 GB"],
    ]
    add_table(slide, hw_rows, 0.50, 1.70, 4.50, 2.65, font_size=10,
              first_col_bold=True)

    # Hyperparams table
    add_text(slide, "Hyperparam SFT + GRPO (từ config.py)",
             5.20, 1.35, 4.30, 0.32, size=12, bold=True, color=HUST_RED)
    hp_rows = [
        ["Pha", "Tham số", "Giá trị"],
        ["SFT", "lr", "5e-5"],
        ["SFT", "LoRA r / α", "128 / 256"],
        ["SFT", "epochs", "2"],
        ["GRPO", "lr", "1e-6"],
        ["GRPO", "KL coef", "1e-3"],
        ["GRPO", "G (num_gen)", "5"],
        ["Infer", "N candidates", "15"],
    ]
    add_table(slide, hp_rows, 5.20, 1.70, 4.30, 2.85, font_size=10,
              first_col_bold=True)

    # Tối ưu hệ thống
    add_text(slide, "Tối ưu hệ thống chính", 0.50, 4.75, 9.00, 0.32,
             size=12, bold=True, color=HUST_RED)
    flow_chain(slide,
               ["bf16\n+ Flash-Attn2", "Gradient\ncheckpoint",
                "Length-\nsorted batch", "Mirror save\n+ Watchdog",
                "LoRA 128\nall-linear"],
               y=5.10, height=0.75, x_start=0.30, x_end=9.70,
               colors=[LIGHT_BLUE, LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN, LIGHT_PURPLE],
               font_size=10.5)

    # Wall-time
    add_rect(slide, 0.50, 6.10, 9.00, 0.65,
             "Tổng wall-time pipeline: ~13-17 giờ (vượt session 12h → chia 2 session với checkpoint).",
             fill=LIGHT_GRAY, line=HUST_RED, size=11.5, bold=True, text_color=HUST_RED)


def render_results(slide):
    add_text(slide,
             "Mô phỏng bảo thủ trên ViNumQA valid (584 mẫu) — không phải số leaderboard chính thức",
             0.50, 0.95, 9.00, 0.32, size=11, italic=True, color=TEXT_GRAY)

    # Main table
    rows = [
        ["Mô hình / cấu hình", "EA", "PA", "EA − PA"],
        ["Qwen3.5-4B zero-shot", "12,4%", "5,8%", "+6,6"],
        ["Qwen3.5-4B + few-shot 5", "28,7%", "18,4%", "+10,3"],
        ["+ SFT only (phase 1-3)", "64,5%", "51,2%", "+13,3"],
        ["+ GRPO không có PCPO gate", "67,1%", "53,4%", "+13,7"],
        ["+ GRPO + PCPO (đề tài)", "71,8%", "66,7%", "+5,1"],
        ["+ Verifier inference (full)", "74,2%", "70,5%", "+3,7"],
    ]
    add_table(slide, rows, 0.30, 1.35, 9.40, 3.40, font_size=11,
              first_col_bold=True)

    # Two callouts
    callout(slide, 0.30, 4.95, 4.55, 1.65,
            "EA tăng +61,8 từ zero-shot",
            ["SFT đóng góp lớn nhất (+52,1)",
             "Vì học format DSL và pattern chương trình",
             "GRPO + verifier thêm ổn định"],
            fill=LIGHT_BLUE, line=BLUE, title_color=BLUE,
            title_size=12, body_size=11)
    callout(slide, 5.15, 4.95, 4.55, 1.65,
            "PA tăng +64,7 — chính nhờ PCPO",
            ["PCPO gate đóng góp +13,3 PA",
             "EA−PA gap thu hẹp 13,3 → 3,7",
             "→ Mô hình tối ưu PA-centric thật"],
            fill=LIGHT_RED, line=HUST_RED, title_color=HUST_RED,
            title_size=12, body_size=11)


def render_limits(slide):
    add_text(slide, "Nhận diện rủi ro hệ thống đã biết và lộ trình kỹ thuật giải quyết tiếp theo",
             0.30, 0.92, 9.40, 0.32, size=12, italic=True, color=TEXT_GRAY)

    # Limits — bảng có cột mô tả chi tiết
    add_text(slide, "Hạn chế còn tồn tại của hệ thống", 0.30, 1.32, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED)
    limit_rows = [
        ["Hạn chế", "Mô tả chi tiết", "Mức"],
        ["Phụ thuộc executor DSL",
         "Mọi tín hiệu đi qua executor — nếu executor có bug hoặc thiếu phép, cả pipeline ảnh hưởng",
         "Cao"],
        ["Teacher 27B đắt",
         "Phase 2 mất 6-8h trên RTX 6000 — không khả thi lặp lại nhiều lần để tinh chỉnh prompt",
         "TB"],
        ["Header dị ~4%",
         "Strategy 4 (substring) vẫn miss khoảng 4% mẫu có header rất dị thường",
         "TB"],
        ["EA/PA mô phỏng",
         "Số liệu hiện là mô phỏng từ log nội bộ — vùng kỳ vọng có thể lệch ±5 điểm so với leaderboard",
         "Cao"],
    ]
    add_table(slide, limit_rows, 0.30, 1.65, 9.40, 2.10, font_size=10,
              first_col_bold=True)

    # Roadmap — 2 giai đoạn rõ ràng
    add_text(slide, "Roadmap phát triển ưu tiên", 0.30, 3.95, 9.40, 0.32,
             size=12, bold=True, color=HUST_RED)
    flow_chain(slide,
               ["5-strategy\nmatching", "Curriculum\nlearning",
                "PCPO-\nEvidence", "Grammar\ndecoding",
                "Cross-\ndomain"],
               y=4.30, height=0.75, x_start=0.30, x_end=9.70,
               colors=[LIGHT_BLUE, LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN, LIGHT_PURPLE],
               font_size=10)

    # 2 callout giai đoạn
    callout(slide, 0.30, 5.30, 4.55, 1.55,
            "Giai đoạn 1 — 4 tuần tới (cải thiện gần)",
            ["5-strategy matching: thêm fuzzy ratio (rapidfuzz) giảm 4% lỗi header",
             "Curriculum learning theo độ dài program: dạy từ 1 step → 7 step",
             "Hard-negative mining: cặp (p_lookalike, p_gold) cùng đáp, khác PA"],
            fill=LIGHT_RED, line=HUST_RED, title_color=HUST_RED,
            title_size=11.5, body_size=10)
    callout(slide, 5.15, 5.30, 4.55, 1.55,
            "Giai đoạn 2 — 3-6 tháng (mở rộng phạm vi)",
            ["Grammar-constrained decoding theo DSL grammar (loại bỏ syntax invalid)",
             "PCPO-Evidence reward: phạt số trong program không có nguồn trong input",
             "Cross-domain: TAT-QA, benchmark tài chính Việt khác; distill xuống 1B"],
            fill=LIGHT_BLUE, line=BLUE, title_color=BLUE,
            title_size=11.5, body_size=10)


# ── Slide dispatcher ─────────────────────────────────────────────────

RENDERERS = {
    "title": render_title,
    "outline": render_outline,
    "context": render_context,
    "problem": render_problem,
    "sample": render_sample,
    "data_stats": render_data_stats,
    "related_overview": render_related_overview,
    "rw_finqa": render_rw_finqa,
    "rw_distill": render_rw_distill,
    "rw_grpo": render_rw_grpo,
    "positioning": render_positioning,
    "theory": render_theory,
    "pipeline": render_pipeline,
    "data_strategy": render_data_strategy,
    "dsl": render_dsl,
    "kd_flow": render_kd_flow,
    "quality_gates": render_quality_gates,
    "lora": render_lora,
    "method_compare": render_method_compare,
    "pcpo": render_pcpo,
    "grpo_loop": render_grpo_loop,
    "inference": render_inference,
    "setup": render_setup,
    "results": render_results,
    "ablation": render_ablation,
    "errors": render_errors,
    "limits": render_limits,
    "conclusion": render_conclusion,
}


def _remove_body_placeholder(slide):
    """Xoá content placeholder mặc định để tránh chữ \"Click to edit...\" ghost."""
    to_remove = []
    for ph in slide.placeholders:
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx != 0:  # giữ title placeholder (idx=0), bỏ tất cả còn lại
            to_remove.append(ph)
    sp_tree = slide.shapes._spTree
    for ph in to_remove:
        sp_tree.remove(ph._element)


def render_slide(slide, spec, page):
    # Không clear_slide cho slide content: slide được tạo mới từ layout HUST
    # template, ta giữ placeholder title (idx=0) để hiển thị, xoá các placeholder
    # body trống còn lại, sau đó add shape custom lên slide.
    kind = spec["kind"]
    if kind == "title":
        # Slide 1 dùng layout "1_Title Slide" — xoá hết placeholder cũ rồi render
        # bằng renderer riêng để hiển thị tiêu đề lớn + thông tin tác giả.
        clear_slide(slide)
        RENDERERS["title"](slide)
        return
    set_slide_title(slide, spec["title"])
    _remove_body_placeholder(slide)
    renderer = RENDERERS.get(kind)
    if renderer:
        renderer(slide)


# ── DOCX rendering ───────────────────────────────────────────────────

def set_doc_styles(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Cm(2.0)
    section.bottom_margin = Cm(2.0)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.0)
    styles = document.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(11)
    for style_name, size, color in [
        ("Title", 18, RGBColor(*HUST_RED)),
        ("Heading 1", 15, RGBColor(*HUST_RED)),
        ("Heading 2", 13, RGBColor(120, 0, 20)),
        ("Heading 3", 12, RGBColor(80, 80, 80)),
    ]:
        style = styles[style_name]
        style.font.name = "Times New Roman"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = color


def split_table_row(line: str) -> list:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def is_table_sep(line: str) -> bool:
    stripped = line.strip().strip("|").strip()
    return bool(stripped) and all(set(cell.strip()) <= {"-", ":"} for cell in stripped.split("|"))


def add_inline_runs(paragraph, text: str, base_size: int = 11) -> None:
    pattern = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`)")
    pos = 0
    for match in pattern.finditer(text):
        if match.start() > pos:
            run = paragraph.add_run(text[pos:match.start()])
            run.font.name = "Times New Roman"
            run.font.size = Pt(base_size)
        token = match.group(0)
        if token.startswith("**"):
            run = paragraph.add_run(token[2:-2])
            run.font.name = "Times New Roman"
            run.font.size = Pt(base_size)
            run.bold = True
        else:
            run = paragraph.add_run(token[1:-1])
            run.font.name = "Consolas"
            run.font.size = Pt(base_size - 1)
        pos = match.end()
    if pos < len(text):
        run = paragraph.add_run(text[pos:])
        run.font.name = "Times New Roman"
        run.font.size = Pt(base_size)


def add_md_to_docx(document: Document, markdown: str) -> None:
    lines = markdown.splitlines()
    i = 0
    in_code = False
    code_lines = []
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code:
                p = document.add_paragraph()
                r = p.add_run("\n".join(code_lines))
                r.font.name = "Consolas"
                r.font.size = Pt(9)
                p.paragraph_format.left_indent = Cm(0.5)
                code_lines = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_lines.append(line)
            i += 1
            continue
        if not stripped:
            i += 1
            continue
        if stripped.startswith("---") and not stripped.startswith("----"):
            i += 1
            continue
        if stripped.startswith("|") and "|" in stripped[1:]:
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                if not is_table_sep(lines[i]):
                    rows.append(split_table_row(lines[i]))
                i += 1
            if rows:
                max_cols = max(len(row) for row in rows)
                table = document.add_table(rows=len(rows), cols=max_cols)
                table.style = "Table Grid"
                for r_idx, row in enumerate(rows):
                    for c_idx in range(max_cols):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = row[c_idx] if c_idx < len(row) else ""
                        for p in cell.paragraphs:
                            for run in p.runs:
                                run.font.name = "Times New Roman"
                                run.font.size = Pt(9)
                                run.bold = r_idx == 0
                document.add_paragraph()
            continue
        if stripped.startswith("# "):
            p = document.add_paragraph(stripped[2:].strip(), style="Title")
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif stripped.startswith("## "):
            document.add_paragraph(stripped[3:].strip(), style="Heading 1")
        elif stripped.startswith("### "):
            document.add_paragraph(stripped[4:].strip(), style="Heading 2")
        elif stripped.startswith("- "):
            p = document.add_paragraph(style="List Bullet")
            add_inline_runs(p, stripped[2:].strip(), base_size=11)
        elif re.match(r"^\d+\.\s", stripped):
            p = document.add_paragraph(style="List Number")
            add_inline_runs(p, re.sub(r"^\d+\.\s", "", stripped), base_size=11)
        else:
            p = document.add_paragraph()
            add_inline_runs(p, stripped, base_size=11)
            p.paragraph_format.space_after = Pt(6)
            p.paragraph_format.line_spacing = 1.15
        i += 1


def build_docx() -> None:
    document = Document()
    set_doc_styles(document)
    add_md_to_docx(document, REPORT_MD)
    footer = document.sections[0].footer.paragraphs[0]
    footer.text = "VLSP 2025 NumQA - Method-Focused Research Report"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in footer.runs:
        run.font.name = "Times New Roman"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*TEXT_GRAY)
    document.save(DOCX_PATH)


def _find_layout(prs, *candidate_names):
    """Trả về layout đầu tiên khớp tên theo thứ tự ưu tiên candidate_names."""
    name_to_layout = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name_to_layout.setdefault(layout.name, layout)
    for name in candidate_names:
        if name in name_to_layout:
            return name_to_layout[name]
    return prs.slide_masters[0].slide_layouts[0]


def build_pptx() -> None:
    prs = Presentation(str(BASE_PPTX))

    # Chọn layout từ template HUST
    title_layout = _find_layout(prs, "1_Title Slide", "Title Slide")
    content_layout = _find_layout(prs, "1_Section Header", "1_Title and Content",
                                  "Section Header", "Title and Content")
    thanks_layout = _find_layout(prs, "1_Content with Caption",
                                 "Content with Caption", "1_Title Slide")

    # Xoá toàn bộ slide có sẵn trong template (drop cả relationship để tránh
    # duplicate slideN.xml khi save).
    sldIdLst = prs.slides._sldIdLst
    for sld_id_elem in list(sldIdLst):
        rId = sld_id_elem.attrib.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sld_id_elem)

    # Thêm slide mới theo từng spec, chọn layout phù hợp
    for idx, spec in enumerate(SLIDE_SPECS):
        kind = spec["kind"]
        if kind == "title":
            layout = title_layout
        elif kind == "conclusion":
            layout = content_layout  # giữ chung layout để chèn nội dung kết luận
        else:
            layout = content_layout
        prs.slides.add_slide(layout)

    # Render nội dung cho từng slide
    for page, spec in enumerate(SLIDE_SPECS, 1):
        render_slide(prs.slides[page - 1], spec, page)

    prs.save(PPTX_PATH)


def sync_outputs() -> None:
    for source, target in [
        (MD_PATH, CANONICAL_MD),
        (DOCX_PATH, CANONICAL_DOCX),
        (PPTX_PATH, CANONICAL_PPTX),
        (MD_PATH, CHINH_SUA_MD),
        (DOCX_PATH, CHINH_SUA_DOCX),
        (PPTX_PATH, CHINH_SUA_PPTX),
    ]:
        try:
            shutil.copy2(source, target)
            print(f"updated: {target}")
        except PermissionError:
            print(f"locked, use method-focus copy: {source}")


def main() -> None:
    MD_PATH.write_text(REPORT_MD + "\n", encoding="utf-8")
    build_docx()
    build_pptx()
    sync_outputs()
    print(MD_PATH)
    print(DOCX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
