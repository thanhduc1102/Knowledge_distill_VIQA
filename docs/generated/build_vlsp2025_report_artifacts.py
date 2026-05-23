from __future__ import annotations

import shutil
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt, RGBColor, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches as PptInches, Pt as PptPt


OUT_DIR = Path(__file__).resolve().parent
ROOT = OUT_DIR.parents[1]
DOCX_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu_chinh_sua.docx"
MD_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu_chinh_sua.md"
PPTX_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_chinh_sua.pptx"
CANONICAL_DOCX_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu.docx"
CANONICAL_MD_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Bao_cao_nghien_cuu.md"
CANONICAL_PPTX_PATH = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST.pptx"
TEMPLATE_PPTX = ROOT / "docs" / "HUST_PPT_template_2022_RED_4x3 (3).pptx"


REPORT_MD = r"""
# Báo cáo nghiên cứu: Knowledge Distillation cho bài toán hỏi đáp suy luận số học tài chính tiếng Việt VLSP 2025

**Định hướng triển khai:** pipeline `data_prep -> teacher -> sft -> grpo -> inference -> evaluate` trong repository `Knowledge_distill_VIQA`, tham khảo bài báo *VLSP 2025 challenge: Numerical Reasoning Question and Answer* và cách tiếp cận của nhóm HUSTUET.

**Ngày biên soạn:** 22/05/2026.

**Phạm vi:** báo cáo tập trung vào bài toán VLSP 2025 Numerical Reasoning QA trong miền tài chính tiếng Việt, dữ liệu ViNumQA/FinQA, kiến trúc Knowledge Distillation, SFT, GRPO/PCPO, suy luận đa đường và đánh giá EA/PA.

**Tuyên bố về số liệu:** các thống kê dữ liệu được trích trực tiếp từ repository và paper. Các bảng kết quả chưa có log chạy thật trong workspace được ghi rõ là **mô phỏng/hypothetical** nhằm phục vụ thiết kế nghiên cứu, không nên trình bày như kết quả chính thức khi chưa chạy lại pipeline. Sau khi hiệu chỉnh, kết quả dự kiến cho pipeline hiện tại được đặt theo hướng bảo thủ: hệ thống hoàn chỉnh chỉ nên dao động khoảng 0,70-0,75 tùy metric và tập đánh giá, tránh các giả định quá lạc quan so với quy mô teacher/student hiện có.

## Tóm tắt điều hành

Bài toán Numerical Reasoning Question Answering trong miền tài chính yêu cầu mô hình đọc đồng thời văn bản trước bảng, bảng số liệu, văn bản sau bảng và câu hỏi tự nhiên để sinh ra hai đầu ra: chương trình tính toán có thể thực thi và đáp án số cuối cùng. Điểm khác biệt quan trọng so với hỏi đáp thông thường là hệ thống không chỉ cần trả lời đúng mà còn phải tạo được chuỗi thao tác toán học minh bạch, có thể kiểm chứng bằng bộ thực thi. Vì vậy, Program Accuracy (PA) trở thành thước đo trọng tâm cho tính đúng đắn của lập luận, trong khi Execution Accuracy (EA) phản ánh tính đúng của đáp án cuối.

Bài báo VLSP 2025 giới thiệu ViNumQA, benchmark đầu tiên cho suy luận số học tài chính tiếng Việt. ViNumQA gồm 4.074 bộ ba câu hỏi - chương trình - đáp án đã được kiểm chứng thủ công, chia thành 2.993 train, 584 valid và 497 test. Tập dữ liệu được xây dựng từ hai nguồn: bản FinQA được dịch và xác thực sang tiếng Việt, và dữ liệu mới trích từ báo cáo tài chính doanh nghiệp Việt Nam giai đoạn 2020-2025. Private test trong repo có 1.625 mẫu không có nhãn chương trình/đáp án.

Hệ thống trong repository định hướng theo chiến lược thắng PA của HUSTUET: mở rộng dữ liệu bằng FinQA tiếng Anh nguyên bản và `program_re`, dùng teacher lớn Qwen3.5-27B hoặc lớn hơn để sinh reasoning trace, fine-tune student Qwen3.5-4B bằng LoRA, sau đó tối ưu bằng GRPO với hàm thưởng PCPO ưu tiên chương trình hợp lệ và chính xác. Ở suy luận, mô hình sinh nhiều ứng viên chương trình, thực thi bằng DSL executor và chọn đáp án bằng majority voting.

Đóng góp cụ thể của hướng triển khai gồm: chuẩn hóa pipeline dữ liệu đa ngôn ngữ; chuyển bảng sang Markdown; thiết kế định dạng đầu ra tiếng Việt gồm phân tích lập luận, chương trình và đáp án; xây dựng DSL tài chính có thể thực thi; áp dụng distillation theo reasoning trace; tối ưu GRPO bằng reward hướng chương trình; đánh giá bằng PA symbolic qua sympy; và hỗ trợ nhiều GPU profile từ P100 16GB đến RTX 6000 96GB.

## 1. Bối cảnh và động lực nghiên cứu

Tài liệu tài chính doanh nghiệp thường tồn tại dưới dạng kết hợp giữa đoạn văn diễn giải và bảng số liệu. Một câu hỏi phân tích tài chính có thể yêu cầu lấy số liệu doanh thu trong một năm, so sánh với năm trước, tính phần trăm tăng trưởng, hoặc kết hợp một giá trị trong bảng với thông tin chỉ xuất hiện trong văn bản. Vì vậy, hệ thống hỏi đáp tài chính cần đồng thời giải quyết ba yêu cầu: hiểu ngôn ngữ tự nhiên, truy xuất đúng số liệu và thực hiện phép tính chính xác.

Trong tiếng Anh, các benchmark như DROP, MathQA, TabFact, FinQA và TAT-QA đã thúc đẩy nghiên cứu về suy luận số học và suy luận trên bảng. Tuy nhiên, tiếng Việt thiếu một benchmark tương đương trong miền tài chính. ViNumQA lấp khoảng trống này bằng cách yêu cầu mô hình sinh chương trình tính toán minh bạch thay vì chỉ sinh câu trả lời. Đây là lựa chọn phù hợp với tài chính, nơi kết quả cần được kiểm toán và truy vết.

Thách thức chính của bài toán nằm ở chỗ mô hình phải học cả ngữ nghĩa tài chính lẫn hình thức chương trình. Một đáp án số đúng nhưng chương trình sai có thể gây rủi ro trong môi trường tài chính vì người dùng không biết mô hình dựa vào logic nào. Ngược lại, chương trình đúng giúp hệ thống có thể tự kiểm tra bằng executor, so sánh với biểu thức symbolic, phát hiện lỗi cú pháp và giải thích đường đi từ input đến output.

## 2. Định nghĩa bài toán

### 2.1 Input

Mỗi mẫu dữ liệu là một đối tượng JSON gồm các trường chính:

| Trường | Ý nghĩa | Vai trò trong mô hình |
|---|---|---|
| `pre_text` | Danh sách đoạn văn đứng trước bảng | Cung cấp bối cảnh, mô tả doanh nghiệp, đơn vị, chỉ tiêu |
| `table` | Bảng tài chính dạng danh sách hai chiều, hàng đầu là header | Nguồn số liệu chính cho phép tính |
| `post_text` | Danh sách đoạn văn sau bảng | Bổ sung diễn giải, nhận xét hoặc số liệu không có trong bảng |
| `id` | Định danh mẫu | Truy vết, merge prediction với gold |
| `qa.question` | Câu hỏi tự nhiên | Xác định mục tiêu tính toán |
| `qa.program` | Chương trình gold, chỉ có ở train/valid/test có nhãn | Nhãn cho SFT, GRPO và PA |
| `qa.exe_ans` | Đáp án thực thi từ chương trình gold | Nhãn cho EA và reward execution |

Trong tập private test, `program` và `exe_ans` được ẩn; hệ thống phải dự đoán cả hai.

### 2.2 Output

Output bắt buộc gồm ba phần bằng tiếng Việt:

```text
**Phân tích lập luận:**
<giải thích từng bước>

**Chương trình tính toán:**
divide(914, 391), multiply(#0, const_100)

**Đáp án cuối cùng:**
233.75959
```

Phần chương trình sử dụng DSL tài chính gồm các hàm `add`, `subtract`, `multiply`, `divide`, `exp`, `greater`, `table_sum`, `table_average`, `table_max`, `table_min`. Kết quả của bước trước được tham chiếu bằng `#0`, `#1`, ...; hằng số dùng dạng `const_100`, `const_1000`, `const_m1`.

### 2.3 Công thức bài toán

Gọi input là `x = (pre_text, table, post_text, question)`. Mô hình cần sinh `y = (r, p, a)`, trong đó `r` là reasoning trace, `p` là chương trình và `a = exec(p, table)` là đáp án sau khi thực thi. Mục tiêu tối ưu không chỉ là cực đại xác suất sinh đúng đáp án mà còn cực đại tính hợp lệ và tương đương symbolic của chương trình:

`p* = argmax_p P_theta(p | x)` với điều kiện `valid(p)=1` và `sym(p) == sym(p_gold)`.

Trong thực tế, chương trình có thể có nhiều biến thể đúng. Do đó, hệ thống đánh giá dùng PA symbolic thay vì so sánh chuỗi tuyệt đối.

## 3. Metric đánh giá

### 3.1 Execution Accuracy (EA)

EA đo tỷ lệ mẫu mà đáp án cuối cùng đúng sau khi thực thi chương trình dự đoán. Với `N` mẫu và `N_correct` mẫu có đáp án đúng:

`EA = N_correct / N * 100%`.

Trong repository, so sánh số thực dùng tolerance tương đối `1e-4`; nếu gold answer bằng 0 thì dùng tolerance tuyệt đối `1e-5`. EA hữu ích vì phản ánh chất lượng trả lời cuối cùng, nhưng có thể đánh giá quá cao năng lực lập luận nếu chương trình sai nhưng tình cờ cho kết quả đúng.

### 3.2 Program Accuracy (PA)

PA đo tỷ lệ chương trình dự đoán tương đương về mặt toán học với chương trình gold. Quy trình gồm token hóa chương trình, thay tham số số học và tham chiếu bảng bằng ký hiệu trừu tượng, dựng biểu thức symbolic và so sánh dạng rút gọn bằng sympy. Ví dụ `add(A, B)` và `add(B, A)` được xem là tương đương vì phép cộng có tính giao hoán.

PA là metric chính vì phản ánh tính kiểm toán của hệ thống. Trong miền tài chính, PA cao có nghĩa là mô hình không chỉ trả lời đúng mà còn chọn đúng logic tính toán.

### 3.3 Hạn chế của metric

PA vẫn có thể false negative khi một câu hỏi có nhiều chương trình khác nhau nhưng đều hợp lệ về mặt ngữ nghĩa tài chính. Ví dụ một tỷ lệ có thể tính trực tiếp từ hai số đã có hoặc gián tiếp qua biến trung gian; nếu chương trình gold chỉ chấp nhận một đường đi, một chương trình khác đúng về thực tế có thể bị đánh rớt. Vì vậy, hướng phát triển nên bổ sung tập chương trình tương đương, verifier dựa trên execution và kiểm tra bằng dữ liệu chứng cứ.

## 4. Phân tích dữ liệu

### 4.1 ViNumQA theo paper VLSP 2025

Bài báo VLSP 2025 mô tả ViNumQA gồm 4.074 bộ ba question-program-answer, chia thành 2.993 train, 584 valid và 497 test. Hai nguồn dữ liệu cân bằng tương đối: 2.069 mẫu Vi Data từ báo cáo tài chính Việt Nam và 2.005 mẫu FinQA-Vi dịch từ FinQA tiếng Anh.

| Split | Vi Data | FinQA-Vi dịch | Tổng |
|---|---:|---:|---:|
| Train | 1.474 | 1.519 | 2.993 |
| Valid | 308 | 276 | 584 |
| Test | 287 | 210 | 497 |
| Tổng | 2.069 | 2.005 | 4.074 |

Theo loại bằng chứng, paper chia câu hỏi thành Table Only, Text Only và Table & Text. Table Only chiếm đa số, nhưng Table & Text là nhóm khó nhất vì cần kết hợp nhiều nguồn.

| Split | Table Only | Table & Text | Text Only | Tổng |
|---|---:|---:|---:|---:|
| Train | 2.213 | 408 | 372 | 2.993 |
| Valid | 441 | 71 | 72 | 584 |
| Test | 365 | 64 | 68 | 497 |

Nhận xét: phân phối nghiêng mạnh về Table Only giúp mô hình học tốt các thao tác trên bảng, nhưng khiến nhóm Table & Text có ít mẫu và dễ trở thành điểm nghẽn khi đánh giá thực tế.

### 4.2 Thống kê dữ liệu hiện có trong workspace

Các thống kê dưới đây được tính trực tiếp từ JSON trong repository.

| Tập dữ liệu | Số mẫu | Số từ câu hỏi TB | pre_text TB | post_text TB | Hàng bảng TB | Cột bảng TB | Bước chương trình TB | Max bước |
|---|---:|---:|---:|---:|---:|---:|---:|---:|
| ViNumQA train | 2.993 | 20,65 | 420,65 | 232,78 | 8,15 | 5,29 | 1,56 | 7 |
| ViNumQA valid | 584 | 20,34 | 399,31 | 206,52 | 8,25 | 5,46 | 1,51 | 6 |
| ViNumQA test | 497 | 20,82 | 432,94 | 172,82 | 8,58 | 5,19 | 1,41 | 5 |
| ViNumQA private test | 1.625 | 19,46 | 388,83 | 116,26 | 9,12 | 6,22 | 0 | 0 |
| FinQA train (Kaggle) | 6.251 | 16,65 | 301,99 | 329,67 | 6,34 | 3,84 | 1,54 | 6 |
| FinQA dev (Kaggle) | 883 | - | - | - | - | - | 1,54 | 5 |
| FinQA test (Kaggle) | 1.147 | - | - | - | - | - | 1,54 | 5 |

Điểm cần lưu ý là thư mục `dataset/finqa_en` trong workspace gốc hiện chưa có file `train/dev/test.json`, nhưng bản Kaggle workspace có đầy đủ `kaggle/kaggle_workspace/data_finqa`. Notebook production sao chép dữ liệu Kaggle vào `dataset/finqa_en`, vì vậy báo cáo triển khai cần phân biệt trạng thái local và trạng thái production.

### 4.3 Phân phối phép toán

Trên ViNumQA train, phép `divide` và `subtract` áp đảo, phản ánh bản chất câu hỏi tài chính thường hỏi tỷ lệ, mức tăng/giảm và chênh lệch.

| Operation | ViNumQA train | ViNumQA valid | ViNumQA test | FinQA train |
|---|---:|---:|---:|---:|
| divide | 1.780 | 336 | 274 | 4.445 |
| subtract | 1.414 | 293 | 242 | 2.739 |
| add | 744 | 126 | 91 | 1.512 |
| multiply | 260 | 29 | 31 | 567 |
| greater | 0 | 0 | 1 | 124 |
| table_average | 126 | 25 | 21 | 95 |
| table_max | 178 | 40 | 22 | 48 |
| table_min | 100 | 21 | 14 | 27 |
| table_sum | 52 | 13 | 6 | 36 |
| exp | 0 | 0 | 0 | 5 |

Hàm `greater` và `exp` ít xuất hiện, dễ dẫn tới thiếu học. Nếu private test có nhiều câu hỏi so sánh hoặc lũy thừa, mô hình có thể suy giảm PA. Cần cân bằng bằng augmentation hoặc sampling có trọng số.

### 4.4 `program_re` và mở rộng dữ liệu

FinQA có trường `program_re` biểu diễn chương trình thay thế cho cùng một đáp án. Trong dữ liệu Kaggle, FinQA train có 6.251 mẫu có `program_re`, trong đó 2.534 mẫu khác chương trình gốc; dev có 360 mẫu khác; test có 493 mẫu khác. Khi bật `use_program_re=True`, pipeline nhân bản mẫu với chương trình thay thế, đưa tổng dữ liệu SFT dự kiến lên khoảng 14.661 mẫu: 2.993 ViNumQA train, 8.281 FinQA all splits và khoảng 3.387 mẫu mở rộng từ `program_re`.

Ý nghĩa của `program_re` là dạy mô hình rằng một câu hỏi tài chính có thể có nhiều đường tính tương đương. Điều này làm giảm overfit vào một chuỗi thao tác duy nhất và phù hợp với PA symbolic. Tuy nhiên, cần kiểm tra `program_re` bằng executor để loại bỏ lỗi do dữ liệu gốc hoặc khác biệt định dạng.

### 4.5 Vấn đề dữ liệu và cách xử lý

| Vấn đề | Biểu hiện | Ảnh hưởng | Cách xử lý trong hệ thống | Hướng cải tiến |
|---|---|---|---|---|
| Dịch thuật/semantic drift | Thuật ngữ tài chính dịch không nhất quán | Mô hình lấy nhầm chỉ tiêu | Dùng FinQA tiếng Anh trực tiếp thay vì dịch lại toàn bộ | Từ điển thuật ngữ tài chính song ngữ và kiểm tra bằng entity alignment |
| Bảng PDF phức tạp | Header nhiều cấp, đơn vị trong ngoặc | Sai row/column retrieval | Chuyển bảng 2D sang Markdown | Table parser giữ cấu trúc multi-index và unit normalization |
| Định dạng số | Dấu phẩy, %, tiền tệ, ngoặc âm | Executor parse sai | `_parse_number`, bỏ ký hiệu tiền, xử lý `%` | Bộ chuẩn hóa số theo locale Việt Nam, phát hiện đơn vị triệu/tỷ |
| Private test thiếu nhãn | Không có `program/exe_ans` | Không thể đánh giá offline | Chỉ inference và lưu prediction | Xây validation nội bộ từ train/valid, cross-validation |
| Mất cân bằng operation | `divide/subtract` nhiều, `greater/exp` ít | Mô hình kém ở phép hiếm | program_re và FinQA tăng đa dạng | Oversampling phép hiếm, synthetic templates |
| Table & Text ít mẫu | Câu hỏi multi-hop thiếu đại diện | PA giảm ở nhóm khó | Distillation reasoning trace | Sinh thêm câu hỏi multi-source và decomposition labels |
| Hallucinated program | Sinh hàm ngoài DSL hoặc sai format | Reward = 0, PA thấp | Regex extraction, `validate_program`, PCPO gating | Grammar-constrained decoding, parser sửa lỗi có kiểm soát |

## 5. Quy trình xử lý và thiết kế dữ liệu trong repository

### 5.1 Pha data preparation

Module `pipeline/data_prep.py` thực hiện bốn nhiệm vụ: tải ViNumQA và FinQA, chuyển bảng sang Markdown, xây prompt theo template tiếng Việt, và sinh ba dạng dữ liệu cho teacher/SFT/GRPO. Hàm `table_to_markdown()` giữ cấu trúc hàng-cột dễ đọc hơn so với list JSON, phù hợp với quan sát trong paper rằng Markdown là biểu diễn bảng hiệu quả cho LLM.

Output chính gồm:

| File | Dạng | Mục đích |
|---|---|---|
| `teacher_input.json` | prompt + ground truth | Đưa vào teacher distillation |
| `sft_train.json` | conversation messages | SFT baseline hoặc fallback khi teacher fail |
| `sft_valid.json` | conversation messages | Validation/inference nội bộ |
| `grpo_train.parquet` | prompt + reward_model ground truth | Huấn luyện GRPO |
| `grpo_valid.parquet` | prompt + reward_model ground truth | Đánh giá reward/GRPO |

### 5.2 Prompt và định dạng output

Prompt trong `src/assets/template.py` yêu cầu mô hình đóng vai chuyên gia phân tích bảng tài chính, trả lời bằng tiếng Việt và tuân thủ ba phần: phân tích lập luận, chương trình tính toán, đáp án cuối cùng. Việc cố định format là thiết yếu vì các module `reward.py`, `teacher_distill.py` và `evaluate.py` đều trích chương trình/đáp án bằng regex.

Phiên bản guided trong `src/assets/think.py` nhúng chương trình gold và đáp án gold vào prompt để teacher tập trung sinh reasoning trace quanh lời giải đúng. Cách này giúp giảm chi phí decode, giảm lỗi chương trình trong distillation và tăng chất lượng SFT rationale.

### 5.3 Teacher distillation

Module `pipeline/teacher_distill.py` dùng teacher local hoặc API để sinh reasoning trace. Validation có bốn mức: `exact_match`, `answer_match`, `program_valid`, `invalid`. Với guided template, `program_valid` vẫn có thể chấp nhận vì chương trình đúng đã được nhúng như hint. Nếu teacher thất bại, pipeline có thể fallback về gold SFT data.

Artifact hiện có cho thấy các run ViNumQA-only đã sinh khoảng 1.998-2.023 mẫu distilled trên 2.993 mẫu train, tương đương 66,8-67,6% usable yield trong snapshot hiện tại. Đây là dấu hiệu cần tối ưu teacher prompt/retry nếu muốn khai thác đủ tập train. Production profile dùng Qwen3.5-27B, batch theo độ dài prompt và checkpoint định kỳ để phù hợp giới hạn thời gian Kaggle.

### 5.4 SFT dataset

SFT dùng format hội thoại:

```json
{
  "messages": [
    {"role": "user", "content": "<prompt>"},
    {"role": "assistant", "content": "<reasoning + program + answer>"}
  ],
  "metadata": {"table": "...", "program": "...", "answer": "..."}
}
```

Trong `train_sft.py`, loss được tính trên assistant response; phần user prompt bị mask bằng `label=-100`. Điều này giúp mô hình học sinh output đúng format và chương trình thay vì học lặp lại input.

### 5.5 GRPO dataset

GRPO dataset lưu prompt và `reward_model.ground_truth`, gồm table, program và answer. Khi mô hình sinh completion, hàm reward sẽ trích chương trình, kiểm tra syntax, thực thi chương trình với table và so sánh đáp án.

Thiết kế này tách dữ liệu huấn luyện sinh văn bản khỏi dữ liệu reward. Đây là ưu điểm vì có thể thay reward mà không phải tái sinh toàn bộ SFT data.

## 6. Nghiên cứu liên quan trước 2026

### 6.1 Dataset và benchmark

| Công trình | Năm | Trọng tâm | Liên hệ với bài toán |
|---|---:|---|---|
| DROP | 2019 | Suy luận rời rạc trên đoạn văn | Nền tảng cho discrete/numerical reasoning |
| MathQA | 2019 | Bài toán số học có operation formalism | Gần với ý tưởng chương trình tính toán |
| TabFact | 2020 | Kiểm chứng phát biểu trên bảng | Nhấn mạnh reasoning trên bảng có cấu trúc |
| FinQA | 2021 | Suy luận số học tài chính text + table, program-based | Nguồn dữ liệu và evaluator chính |
| TAT-QA | 2021 | Hỏi đáp tài chính kết hợp bảng và văn bản | Làm rõ độ khó của Table & Text |
| NumGLUE | 2022 | Bộ benchmark suy luận số học đa nhiệm | Cung cấp góc nhìn tổng quát về numerical reasoning |
| MultiHiertt | 2022 | Suy luận trên bảng phân cấp và văn bản | Gần với bảng tài chính nhiều cấp |
| CT2C-QA | 2024 | QA đa phương thức text-table-chart tiếng Trung | Gợi ý mở rộng sang biểu đồ/báo cáo đa modal |
| ViNumQA/VLSP 2025 | 2025 | Benchmark suy luận số học tài chính tiếng Việt | Đối tượng nghiên cứu trực tiếp |

### 6.2 Kỹ thuật suy luận và sinh chương trình

Các hướng chính gồm semantic parsing sinh chương trình, Chain-of-Thought (CoT), Program-of-Thought/PAL, self-consistency và agent decomposition. Semantic parsing phù hợp với PA vì output có thể thực thi, nhưng cần dữ liệu gán nhãn chương trình. CoT cải thiện giải thích nhưng nếu không ràng buộc format có thể hallucinate. Program-of-Thought tách lập luận ngôn ngữ và tính toán bằng code, nhưng trong VLSP chương trình phải thuộc DSL cố định. Self-consistency sinh nhiều lời giải và vote giúp tăng EA, đặc biệt với câu hỏi khó. Agent decomposition hiệu quả cho Table & Text nhưng chi phí inference cao.

### 6.3 Distillation và fine-tuning

Knowledge distillation từ teacher lớn sang student nhỏ giúp giữ năng lực reasoning trong giới hạn tham số. Có ba biến thể đáng chú ý: distill logits, distill output label và distill reasoning trace. Với bài toán này, distill reasoning trace phù hợp nhất vì student cần học cả quá trình chọn số liệu và ánh xạ sang hàm DSL. LoRA/QLoRA giúp fine-tune mô hình 4B-9B trên GPU giới hạn bằng cách cập nhật ma trận low-rank thay vì toàn bộ trọng số.

### 6.4 Reinforcement learning cho reasoning

Trước 2026, các hướng RL nổi bật gồm PPO/RLHF, DPO, RLVR và GRPO. PPO cần value model và reward model, tốn bộ nhớ. DPO tối ưu trực tiếp cặp chosen/rejected, đơn giản nhưng cần preference pairs chất lượng. RLVR dùng reward có thể kiểm chứng bằng thực thi hoặc đáp án, phù hợp toán học nhưng nếu chỉ thưởng đáp án sẽ bỏ qua PA. GRPO chuẩn hóa reward trong nhóm completion cho cùng prompt, không cần value model, rất phù hợp với bài toán sinh nhiều chương trình và có reward thực thi được.

## 7. Kiến trúc hệ thống hiện tại

Pipeline gồm sáu pha:

| Pha | Module | Input | Output | Mục tiêu |
|---|---|---|---|---|
| 1. Data prep | `pipeline/data_prep.py` | ViNumQA, FinQA | teacher/SFT/GRPO data | Chuẩn hóa dữ liệu và prompt |
| 2. Teacher | `pipeline/teacher_distill.py` | teacher_input | distilled_sft | Sinh reasoning trace chất lượng cao |
| 3. SFT | `pipeline/train_sft.py` | distilled_sft/sft_train | SFT LoRA checkpoint | Học format, program và reasoning |
| 4. GRPO | `pipeline/train_grpo.py` | grpo parquet | GRPO checkpoint | Tối ưu reward hướng chương trình |
| 5. Inference | `pipeline/inference.py` | valid/test prompt | predictions | Sinh N ứng viên, vote đáp án |
| 6. Evaluate | `pipeline/evaluate.py` | predictions + dataset | EA/PA | Đánh giá execution và symbolic program |

### 7.1 Mô hình

Profile production hiện dùng Qwen3.5-27B làm teacher và Qwen3.5-4B làm student trên RTX 6000 96GB. Profile P100 dùng teacher 4B quantized 4-bit và student 0.8B để kiểm thử. LoRA được dùng ở SFT và GRPO để giảm VRAM; Flash Attention bật trên RTX/A100/H100 và tắt trên P100.

| Profile | Teacher | Student | dtype | Flash Attention | Inference candidates |
|---|---|---|---|---|---:|
| p100_16gb | Qwen3.5-4B 4bit | Qwen3.5-0.8B | fp16 | Không | 5 |
| rtx6000_96gb | Qwen3.5-27B | Qwen3.5-4B | bf16 | Có | 15 |
| a100_80gb | Qwen3.5-27B | Qwen3.5-9B | bf16 | Có | 15 |

### 7.2 Financial DSL và executor

`pipeline/program_executor.py` định nghĩa executor cho các phép toán tài chính. Executor parse số, xử lý hằng số, tham chiếu bước trước và một số định dạng tài chính như ký hiệu tiền, dấu phẩy, phần trăm. Output được format tối đa 5 chữ số thập phân, bỏ số 0 thừa và giữ số nguyên ở dạng `5.0`.

Điểm mạnh của DSL là tính kiểm soát: mô hình không được sinh code Python tùy ý, giảm rủi ro bảo mật và dễ đánh giá PA. Điểm yếu là DSL hiện chưa biểu diễn tốt các thao tác như chọn ô theo hàng-cột cụ thể, chuyển đổi đơn vị phức tạp, hoặc so sánh nhiều điều kiện.

## 8. Phương pháp áp dụng

### 8.1 Knowledge Distillation

Teacher lớn sinh reasoning trace có cấu trúc, student học lại trace này bằng SFT. Lý do chọn KD là bài toán cần năng lực reasoning cao nhưng subtask constrained giới hạn tham số và cấm API lúc inference. Teacher chỉ dùng trong huấn luyện; student nhỏ tự chạy khi dự đoán.

So sánh kỹ thuật distillation:

| Kỹ thuật | Ưu điểm | Nhược điểm | Mức phù hợp |
|---|---|---|---|
| Logits KD | Truyền phân phối xác suất mịn | Cần truy cập logits teacher, tốn lưu trữ | Trung bình |
| Output KD | Dễ triển khai, chỉ cần answer/program | Ít học được reasoning | Trung bình |
| Rationale KD | Student học cách chọn số liệu và lập luận | Phụ thuộc chất lượng trace | Cao |
| Program-guided rationale KD | Giảm hallucination chương trình, tăng format compliance | Có nguy cơ teacher chỉ diễn giải lại gold, giảm đa dạng | Rất cao |
| Self-distillation | Không cần teacher quá lớn | Dễ khuếch đại lỗi của student | Trung bình |

Hệ thống chọn rationale KD có guided template vì trực tiếp phục vụ PA: student học cách ánh xạ câu hỏi và dữ liệu sang DSL.

### 8.2 Supervised Fine-Tuning với LoRA

SFT tối ưu log-likelihood của assistant output:

`L_SFT = - sum_t log P_theta(y_t | y_<t, x)`.

LoRA cập nhật ma trận trọng số bằng thành phần low-rank `Delta W = B A`, trong đó `rank r` nhỏ hơn nhiều so với kích thước ma trận gốc. Repository dùng `r=64` cho P100 và `r=128` cho RTX 6000, target `all-linear`, dropout 0,05, gradient checkpointing và cosine learning rate.

So sánh fine-tuning:

| Phương pháp | VRAM | Chất lượng kỳ vọng | Rủi ro | Lý do chọn/không chọn |
|---|---|---|---|---|
| Full fine-tune | Rất cao | Cao nếu đủ dữ liệu | Overfit, tốn checkpoint | Không phù hợp P100/RTX nếu model lớn |
| LoRA | Thấp-vừa | Cao cho task hẹp | Rank quá thấp giảm năng lực | Được chọn vì cân bằng chất lượng/chi phí |
| QLoRA | Thấp | Gần LoRA nếu quant tốt | Lỗi quantization, chậm hơn | Hữu ích cho P100 hoặc student lớn |
| Prefix/Prompt tuning | Rất thấp | Kém hơn với program format phức tạp | Khó học DSL dài | Không ưu tiên |
| Adapter truyền thống | Vừa | Ổn định | Thêm module inference | Có thể thử sau LoRA |

### 8.3 GRPO với reward PCPO

GRPO sinh `N` completion cho cùng prompt, tính reward từng completion, chuẩn hóa reward trong nhóm và cập nhật policy theo advantage tương đối:

`A_i = (R_i - mean(R)) / (std(R) + epsilon)`.

Hàm thưởng PCPO trong `pipeline/reward.py`:

`R(p, x) = R_valid * (0.7 + 0.2 * R_exec + 0.1 * R_bonus)`.

Trong đó `R_valid` là cổng tuyệt đối: chương trình sai cú pháp nhận 0. `R_exec` thưởng đáp án đúng sau thực thi. `R_bonus` thưởng chương trình ngắn hơn hoặc bằng gold để khuyến khích súc tích. Cách thiết kế này bám mục tiêu tài chính: chương trình hợp lệ và kiểm toán được quan trọng hơn việc tình cờ trả lời đúng.

So sánh RL/fine alignment:

| Phương pháp | Cần preference pair | Cần value model | Reward kiểm chứng được | Ưu điểm | Nhược điểm | Phù hợp PA |
|---|---|---|---|---|---|---|
| DPO | Có | Không | Gián tiếp | Đơn giản, ổn định | Cần chosen/rejected chất lượng; khó tối ưu execution trực tiếp | Trung bình |
| PPO/RLHF | Không bắt buộc | Có | Có thể | Linh hoạt | Tốn VRAM, nhiều hyperparameter | Trung bình-cao |
| GPO/Group preference | Có/nhóm | Không hoặc ít | Gián tiếp | Khai thác so sánh nhóm | Không chuẩn hóa execution rõ như GRPO | Trung bình |
| RLVR execution-only | Không | Tùy | Có | Tăng EA mạnh | Dễ bỏ qua chương trình đúng | Trung bình |
| GRPO + PCPO | Không | Không | Có | Tối ưu nhóm completion, ít VRAM, hợp với DSL | Phụ thuộc reward; nếu reward sai sẽ học lệch | Cao |

GRPO được chọn vì không cần value model, phù hợp GPU giới hạn, tận dụng được nhiều completion cho cùng prompt và kết hợp tự nhiên với executor/reward của bài toán.

### 8.4 Multi-path inference và majority voting

Ở inference, mô hình sinh nhiều chương trình với temperature 0,7 và top-p 0,95. Mỗi chương trình hợp lệ được thực thi, đáp án được format rồi vote theo tần suất. Nếu nhiều completion khác nhau cho cùng đáp án, hệ thống chọn đáp án phổ biến nhất và chương trình đầu tiên tương ứng.

Kỹ thuật này thường tăng EA vì giảm lỗi ngẫu nhiên, nhưng có thể tăng PA ít hơn nếu nhiều chương trình khác nhau cùng cho một đáp án. Do đó, voting nên kết hợp với tiêu chí chọn chương trình có reward cao nhất hoặc symbolic closest nếu có verifier phụ.

## 9. So sánh phương pháp và kết quả

### 9.1 Mốc tham chiếu từ VLSP 2025 paper

Bảng dưới đây chỉ dùng làm mốc tham chiếu từ paper, không phải kết quả của pipeline hiện tại. Để tránh suy diễn quá mức, báo cáo ưu tiên so sánh theo PA vì đây là metric chính của hướng nghiên cứu kiểm toán chương trình.

| Subtask | Team | PA paper (%) | Vai trò tham chiếu |
|---|---|---:|---|
| 1 | HUSTUET | 76,63 | Mốc trên về PA cho hướng KD + GRPO với teacher rất lớn |
| 1 | Vietnam Finance | 75,00 | Mốc cạnh tranh khi tối ưu execution và majority voting |
| 1 | Innovation-LLM | 69,82 | Mốc agent inference-only trong constrained setting |
| 1 | UIT_BlackCoffee | 69,67 | Mốc Qwen quantized + Markdown table |
| 2 | Innovation-LLM | 74,07 | Mốc agent decomposition trong unconstrained setting |

Nhận xét quan trọng là các mốc PA của nhóm top nằm chủ yếu quanh 0,70-0,76. Pipeline trong repository dùng teacher 27B và student 4B, nhỏ hơn cấu hình HUSTUET 235B -> 8B trong paper, đồng thời artifact distillation hiện có mới đạt khoảng 67% mẫu usable ở snapshot Vi-only. Vì vậy, kết quả mô phỏng của pipeline hiện tại cần đặt thấp hơn mốc thắng cuộc và chỉ nên xem vùng 0,70-0,75 là hợp lý.

### 9.2 Kết quả mô phỏng cho pipeline hiện tại

Bảng sau là **mô phỏng có kiểm soát** sau khi hiệu chỉnh bảo thủ, dựa trên bốn căn cứ: student hiện tại chỉ 4B, teacher production là 27B thay vì 235B, dữ liệu FinQA local cần đồng bộ từ Kaggle, và distilled SFT snapshot hiện mới có khoảng 2.0K/2.993 mẫu usable ở run Vi-only. Bảng dùng để định hướng ablation; cần chạy lại `pipeline.run` để có kết quả chính thức.

| Thiết lập | EA valid (%) | PA valid (%) | Valid program (%) | Ghi chú |
|---|---:|---:|---:|---|
| Zero-shot Qwen3.5-4B | 52,0 | 35,0 | 55,0 | Dễ sai format và hallucinate hàm |
| Prompt + Markdown + executor voting | 60,5 | 46,8 | 65,0 | Markdown giúp truy xuất bảng tốt hơn nhưng chưa học DSL đủ sâu |
| SFT gold ViNumQA-only | 68,7 | 60,8 | 70,5 | Học format và DSL, nhưng còn thiếu đa dạng chương trình |
| SFT ViNumQA + FinQA English | 70,8 | 63,9 | 72,1 | Transfer reasoning đa ngôn ngữ, tăng vừa phải |
| SFT + `program_re` | 71,6 | 65,4 | 72,8 | Tăng đa dạng chương trình tương đương |
| KD rationale 27B -> 4B | 72,8 | 68,6 | 74,0 | Reasoning trace tăng PA nhưng còn phụ thuộc yield teacher |
| KD + GRPO/PCPO | 73,6 | 70,2 | 74,6 | Reward ưu tiên chương trình hợp lệ và đáp án đúng |
| KD + GRPO/PCPO + vote N=15 | 74,4 | 70,8 | 74,6 | Voting tăng EA vừa phải, PA tăng nhẹ |

Kịch bản hiệu chỉnh này hợp lý hơn vì không giả định student 4B có thể chạm ngay mức của hệ thống thắng cuộc dùng teacher lớn hơn nhiều. Vùng kỳ vọng tốt nhất cho pipeline hiện tại là EA khoảng 0,73-0,75 và PA khoảng 0,70-0,71 trên validation nội bộ; nếu dữ liệu FinQA được đồng bộ đầy đủ, teacher yield tăng và reward được tối ưu, PA có thể tiến gần 0,72-0,74 nhưng chưa nên báo cáo vượt mốc này khi chưa có log thực nghiệm.

### 9.3 Ablation reward mô phỏng

| Reward | EA (%) | PA (%) | Lỗi chính |
|---|---:|---:|---|
| Execution-only `R=R_exec` | 74,2 | 67,3 | Đáp án đúng nhưng program không khớp gold |
| Format-only `R=R_valid` | 66,1 | 61,5 | Chương trình hợp lệ nhưng sai số liệu |
| `0.5 valid + 0.4 exec + 0.1 bonus` | 73,8 | 68,8 | Nghiêng EA, PA chưa tối ưu |
| PCPO `0.7/0.2/0.1` | 73,6 | 70,2 | Cân bằng tốt hơn cho PA nhưng không làm EA tăng đột biến |
| PCPO + evidence reward đề xuất | 74,1 | 71,4 | Cần thêm extractor chứng cứ, chỉ là hướng kỳ vọng |

PCPO được chọn không phải vì tạo ra mức tăng quá lớn, mà vì ổn định hơn với mục tiêu PA: nó phạt sai format, vẫn thưởng execution và tránh chương trình quá dài. Nếu mục tiêu là ứng dụng thực tế thiên về đáp án, có thể tăng `beta`; nếu mục tiêu là leaderboard PA và kiểm toán, giữ hoặc tăng `alpha` hợp lý hơn.

### 9.4 Phân tích theo loại câu hỏi

Paper cho thấy Table Only dễ nhất, Table & Text khó nhất. Với pipeline hiện tại, không nên giả định lặp lại mức Table Only rất cao của hệ thống top paper; thay vào đó cần đặt kỳ vọng theo năng lực student 4B và tình trạng dữ liệu hiện có.

| Nhóm câu hỏi | PA dự kiến bảo thủ (%) | Điểm mạnh của pipeline | Điểm yếu còn lại |
|---|---:|---|---|
| Table Only | 73,5-75,0 | Markdown table, executor rõ ràng | Header nhiều cấp, đơn vị trong bảng |
| Text Only | 64,0-67,0 | Teacher trace giúp chọn số trong văn bản | Nhiều số gần nhau, đại từ/ẩn dụ tài chính |
| Table & Text | 52,0-56,0 | Có thể dùng reasoning trace | Cần multi-hop, hiện chưa có planner tách bước |
| Tổng hợp | 70,0-71,5 | KD + PCPO hỗ trợ PA | Phụ thuộc chất lượng distillation và reward |

## 10. Phân tích lỗi từ input đến output

### 10.1 Lỗi truy xuất số liệu

Mô hình có thể chọn nhầm cột năm, nhầm chỉ tiêu gần tên, hoặc bỏ qua đơn vị. Ví dụ câu hỏi hỏi "năm 2023 gấp bao nhiêu lần năm 2022" nhưng mô hình lấy gross profit thay vì net revenue. Lỗi này thường làm cả EA và PA sai.

### 10.2 Lỗi lập trình DSL

Các lỗi phổ biến gồm sai thứ tự đối số `divide`, dùng `subtract(a,b)` thay vì `subtract(b,a)`, tham chiếu `#1` khi mới có một bước, hoặc sinh hàm ngoài DSL. PCPO xử lý bằng `R_valid` gating, nhưng inference vẫn cần grammar-constrained decoding để giảm lỗi từ đầu.

### 10.3 Lỗi chương trình đúng đáp án nhưng sai PA

Một chương trình có thể ra đáp án đúng nhưng không tương đương gold do dùng đường tính khác hoặc do PA chưa nhận ra tương đương ngữ nghĩa. Đây là lý do nên mở rộng gold bằng `program_re` và tạo verifier đa chương trình.

### 10.4 Lỗi hallucination reasoning

Phần phân tích lập luận có thể nói đến số liệu không có trong input, dù chương trình cuối đúng. Vì reward hiện tập trung vào program và answer, hallucination trong reasoning chưa bị phạt. Hướng cải tiến là evidence-grounded reward: mỗi số trong reasoning/program phải map được về cell hoặc text span.

## 11. Định hướng phát triển

### 11.1 Tối ưu reward

Reward nên mở rộng từ PCPO sang PCPO-Evidence: ngoài valid, execution và conciseness, thêm điểm cho việc trích đúng chứng cứ. Một công thức đề xuất:

`R = R_valid * (0.55 R_form + 0.20 R_exec + 0.15 R_pa_proxy + 0.10 R_evidence)`.

Trong đó `R_pa_proxy` so khớp symbolic với nhiều chương trình tương đương; `R_evidence` kiểm tra số trong chương trình có tồn tại trong bảng/văn bản sau chuẩn hóa đơn vị.

### 11.2 Sinh và xây dựng thêm dữ liệu

Cần tăng dữ liệu Table & Text bằng cách sinh câu hỏi multi-hop có kiểm soát từ báo cáo Việt Nam. Quy trình nên gồm: trích bảng bằng parser, nhận diện đơn vị, sinh câu hỏi bằng LLM, sinh program bằng template hoặc solver, thực thi để kiểm tra, sau đó human review các mẫu khó. Ngoài ra nên tạo counterfactual samples thay đổi năm/chỉ tiêu để giảm học thuộc.

### 11.3 Tăng cường reasoning và giảm hallucination

Nên đưa agent decomposition vào trước chương trình: tách câu hỏi thành subgoal, map từng subgoal vào chứng cứ, rồi sinh DSL. Ở training, teacher có thể sinh bốn trường: intent, evidence, derivation và program. Ở inference, verifier kiểm tra từng số trong program với evidence map trước khi vote.

### 11.4 Cải thiện executor và PA

Executor nên hỗ trợ chọn cell theo `(row, column)`, chuẩn hóa đơn vị triệu/tỷ, số âm trong ngoặc, bảng nhiều header và alias chỉ tiêu. PA nên chấp nhận tập chương trình tương đương sinh bởi rewrite rules: giao hoán cộng/nhân, phân phối, đảo tỷ lệ hợp lệ và rút gọn phép nhân/chia với 100.

### 11.5 Định hướng tốt nhất đề xuất

Hướng có triển vọng nhất là kết hợp KD + GRPO/PCPO với agentic verifier nhẹ. Teacher tạo reasoning trace và evidence map; student học bằng LoRA; GRPO tối ưu chương trình bằng reward thực thi; inference sinh N chương trình nhưng chọn bằng reward verifier thay vì chỉ vote đáp án. Cách này giữ ưu điểm PA của HUSTUET, đồng thời mượn khả năng decomposition của Innovation-LLM cho nhóm Table & Text.

## 12. Đóng góp cụ thể

1. Chuẩn hóa bài toán VLSP 2025 thành pipeline nghiên cứu end-to-end từ dữ liệu đến đánh giá.
2. Phân tích dữ liệu ViNumQA/FinQA hiện có, gồm số mẫu, độ dài, bảng, bước chương trình và phân phối operation.
3. Thiết kế dữ liệu SFT/GRPO/teacher rõ ràng, có thể tái lập và mở rộng bằng `program_re`.
4. Áp dụng distillation reasoning trace để chuyển năng lực teacher lớn sang student nhỏ tự chạy offline.
5. Xây dựng reward PCPO đặt program validity và PA-oriented reasoning lên trước đáp án cuối.
6. Tích hợp executor DSL và symbolic PA để đánh giá minh bạch.
7. Đề xuất roadmap cải thiện reward, dữ liệu multi-source, verifier chứng cứ và grammar-constrained decoding.

## 13. Kết luận

VLSP 2025 Numerical Reasoning QA là bài toán quan trọng cho tài chính tiếng Việt vì yêu cầu mô hình vừa hiểu dữ liệu bảng/văn bản, vừa thực hiện tính toán chính xác, vừa giải thích bằng chương trình kiểm chứng được. Paper VLSP 2025 cho thấy hướng Knowledge Distillation + SFT + GRPO đạt PA tốt nhất trong môi trường constrained, còn agent inference-only đạt EA cao nhất khi được tự do về suy luận.

Hệ thống hiện tại trong repository đi theo hướng phù hợp cho nghiên cứu và triển khai: dùng Qwen3.5 teacher/student, LoRA, guided reasoning distillation, PCPO reward và majority voting. Điểm cần hoàn thiện nằm ở dữ liệu FinQA local, yield distillation, reward evidence-aware, nhóm Table & Text và kiểm soát hallucination. Với các cải tiến này, pipeline có thể trở thành một hệ thống tài chính tiếng Việt vừa chính xác vừa kiểm toán được.

## Phụ lục A. Lệnh chạy chính

```bash
python -m pipeline.run --gpu-profile p100_16gb --phases all --max-samples 50
python -m pipeline.run --gpu-profile rtx6000_96gb --phases all
```

## Phụ lục B. File quan trọng trong repository

| File | Vai trò |
|---|---|
| `pipeline/config.py` | Dataclass config và GPU profiles |
| `pipeline/data_prep.py` | Load/merge dữ liệu, Markdown table, build SFT/GRPO/teacher data |
| `pipeline/teacher_distill.py` | Teacher generation, validation tiers, distilled SFT |
| `pipeline/train_sft.py` | LoRA SFT, masking prompt tokens |
| `pipeline/train_grpo.py` | TRL GRPO hoặc manual REINFORCE fallback |
| `pipeline/reward.py` | PCPO reward |
| `pipeline/program_executor.py` | DSL executor và format answer |
| `pipeline/evaluate.py` | EA/PA, symbolic comparison |
| `src/assets/template.py` | Prompt tiếng Việt chuẩn |
| `src/assets/think.py` | Guided prompt cho teacher reasoning |

## Tài liệu tham khảo chọn lọc

- Amini et al. (2019). MathQA: Towards Interpretable Math Word Problem Solving with Operation-based Formalisms.
- Chen et al. (2020). TabFact: A Large-scale Dataset for Table-based Fact Verification.
- Chen et al. (2021). FinQA: A Dataset of Numerical Reasoning over Financial Data.
- Dua et al. (2019). DROP: A Reading Comprehension Benchmark Requiring Discrete Reasoning Over Paragraphs.
- Hu et al. (2021). LoRA: Low-Rank Adaptation of Large Language Models.
- Robertson and Zaragoza (2009). The Probabilistic Relevance Framework: BM25 and Beyond.
- Shao et al. (2024). DeepSeekMath: Pushing the Limits of Mathematical Reasoning in Open Language Models.
- Wang et al. (2022). Self-Consistency Improves Chain of Thought Reasoning in Language Models.
- Wei et al. (2022). Chain-of-Thought Prompting Elicits Reasoning in Large Language Models.
- Yang et al. (2025). Qwen3 Technical Report.
- Zhao et al. (2022). MultiHiertt: Numerical Reasoning over Multi Hierarchical Tabular and Textual Data.
- Zhao et al. (2024). CT2C-QA: Multimodal Question Answering over Chinese Text, Table and Chart.
- Zhu et al. (2021). TAT-QA: A Question Answering Benchmark on a Hybrid of Tabular and Textual Content in Finance.
- Le Ngoc Toan et al. (2025). VLSP 2025 Challenge: Numerical Reasoning Question and Answer.
""".strip()


SLIDES = [
    {
        "title": "Knowledge Distillation cho VLSP 2025 NumQA",
        "subtitle": "Báo cáo nghiên cứu và định hướng triển khai hệ thống hỏi đáp suy luận số học tài chính tiếng Việt",
        "bullets": ["Dữ liệu ViNumQA/FinQA", "KD + SFT + GRPO/PCPO", "Đánh giá EA/PA và định hướng cải tiến"],
    },
    {
        "title": "1. Bài toán nghiên cứu",
        "bullets": [
            "Input: pre_text, table, post_text và câu hỏi tài chính tiếng Việt.",
            "Output kép: chương trình tính toán có thể thực thi và đáp án số cuối cùng.",
            "Mục tiêu không chỉ là trả lời đúng mà còn sinh reasoning path minh bạch, kiểm toán được.",
            "Metric chính: Program Accuracy (PA); metric phụ: Execution Accuracy (EA).",
        ],
    },
    {
        "title": "2. ViNumQA theo paper VLSP 2025",
        "table": [
            ["Split", "Vi Data", "FinQA-Vi", "Tổng"],
            ["Train", "1.474", "1.519", "2.993"],
            ["Valid", "308", "276", "584"],
            ["Test", "287", "210", "497"],
            ["Tổng", "2.069", "2.005", "4.074"],
        ],
        "bullets": [
            "Dữ liệu gồm hai nguồn: báo cáo tài chính Việt Nam 2020-2025 và FinQA dịch/xác thực.",
            "Table Only chiếm đa số; Table & Text là nhóm khó nhất vì cần multi-hop reasoning.",
        ],
    },
    {
        "title": "3. Thống kê dữ liệu trong repo",
        "table": [
            ["Tập", "Mẫu", "Q words", "Rows", "Cols", "Steps"],
            ["Vi train", "2.993", "20,65", "8,15", "5,29", "1,56"],
            ["Vi valid", "584", "20,34", "8,25", "5,46", "1,51"],
            ["Vi test", "497", "20,82", "8,58", "5,19", "1,41"],
            ["Private", "1.625", "19,46", "9,12", "6,22", "-"],
            ["FinQA train", "6.251", "16,65", "6,34", "3,84", "1,54"],
        ],
        "bullets": ["Operation phổ biến nhất: divide và subtract.", "FinQA có program_re giúp mở rộng dữ liệu lên khoảng 14,6K mẫu."],
    },
    {
        "title": "4. Vấn đề dữ liệu cần giải quyết",
        "bullets": [
            "Translation drift và thuật ngữ tài chính không nhất quán.",
            "Bảng PDF có header phức tạp, đơn vị trong ngoặc và định dạng số đa dạng.",
            "Mất cân bằng operation: divide/subtract áp đảo; greater/exp rất ít.",
            "Table & Text ít mẫu nhưng là nhóm quyết định năng lực reasoning thực tế.",
            "Private test không có program/exe_ans nên cần validation nội bộ đáng tin cậy.",
        ],
    },
    {
        "title": "5. Pipeline dữ liệu",
        "bullets": [
            "Load ViNumQA + FinQA, merge trực tiếp English/Vietnamese thay vì dịch lại.",
            "Chuyển table 2D sang Markdown để LLM đọc hàng-cột ổn định hơn.",
            "Sinh ba format: teacher_input.json, sft_train.json, grpo_train.parquet.",
            "Bật use_program_re để thêm chương trình tương đương từ FinQA.",
            "Guided teacher prompt nhúng gold program/answer để sinh reasoning trace sạch hơn.",
        ],
    },
    {
        "title": "6. Nghiên cứu liên quan",
        "bullets": [
            "Benchmark: DROP, MathQA, TabFact, FinQA, TAT-QA, MultiHiertt, CT2C-QA, ViNumQA.",
            "Reasoning: semantic parsing, CoT, Program-of-Thought/PAL, self-consistency, agent decomposition.",
            "Training: LoRA/QLoRA, rationale distillation, self-distillation.",
            "Alignment/RL: PPO, DPO, RLVR, GRPO; GRPO phù hợp vì không cần value model và có reward kiểm chứng được.",
        ],
    },
    {
        "title": "7. Kiến trúc hệ thống",
        "bullets": [
            "Data prep -> Teacher distill -> SFT -> GRPO/PCPO -> Inference voting -> Evaluate.",
            "Teacher production: Qwen3.5-27B; student: Qwen3.5-4B.",
            "LoRA giúp huấn luyện student trong giới hạn VRAM.",
            "Executor DSL đảm bảo chương trình sinh ra có thể chạy và kiểm tra PA symbolic.",
        ],
    },
    {
        "title": "8. DSL và metric",
        "table": [
            ["Nhóm", "Hàm"],
            ["Số học", "add, subtract, multiply, divide, exp"],
            ["So sánh", "greater"],
            ["Bảng", "table_sum, table_average, table_max, table_min"],
            ["Tham chiếu", "#0, #1, const_100, const_m1"],
        ],
        "bullets": ["EA đo đúng đáp án cuối; PA đo tương đương symbolic của chương trình.", "PA phù hợp hơn cho tài chính vì ưu tiên kiểm toán logic."],
    },
    {
        "title": "9. Knowledge Distillation",
        "bullets": [
            "Teacher lớn sinh reasoning trace theo format tiếng Việt.",
            "Validation tiers: exact_match > answer_match > program_valid > invalid.",
            "Student học không chỉ program mà cả cách xác định số liệu và ánh xạ sang DSL.",
            "Snapshot hiện có: khoảng 2.023/2.993 mẫu ViNumQA distilled usable trong một run Vi-only.",
        ],
    },
    {
        "title": "10. SFT với LoRA",
        "bullets": [
            "Loss chỉ tính trên assistant response; user prompt bị mask bằng label=-100.",
            "LoRA cập nhật Delta W = B A, giảm số tham số trainable và VRAM.",
            "Profile RTX 6000: LoRA r=128, alpha=256, bf16, flash attention.",
            "Profile P100: LoRA r=64, fp16, không flash attention.",
        ],
    },
    {
        "title": "11. GRPO với PCPO reward",
        "bullets": [
            "GRPO sinh nhiều completion cho cùng prompt và chuẩn hóa reward trong nhóm.",
            "PCPO: R = R_valid * (0.7 + 0.2 R_exec + 0.1 R_bonus).",
            "R_valid là cổng tuyệt đối: sai cú pháp nhận reward 0.",
            "Lý do chọn GRPO: không cần value model, nhẹ VRAM, hợp với reward thực thi DSL.",
        ],
    },
    {
        "title": "12. Suy luận đa đường",
        "bullets": [
            "Sinh N chương trình với temperature/top-p để tạo đa dạng ứng viên.",
            "Validate và execute từng program bằng DSL executor.",
            "Majority vote trên đáp án thực thi, confidence = tỷ lệ đồng thuận.",
            "Voting tăng EA rõ hơn PA; nên kết hợp verifier chọn chương trình tốt nhất.",
        ],
    },
    {
        "title": "13. So sánh phương pháp",
        "table": [
            ["Nhóm", "Tốt nhất cho", "Rủi ro"],
            ["KD rationale", "PA, format, reasoning", "Phụ thuộc teacher"],
            ["LoRA SFT", "Chi phí thấp", "Rank thấp gây thiếu học"],
            ["DPO", "Preference ổn định", "Cần cặp chosen/rejected"],
            ["RLVR", "EA", "Có thể bỏ qua PA"],
            ["GRPO/PCPO", "PA + reward thực thi", "Cần reward chuẩn"],
        ],
    },
    {
        "title": "14. Kết quả mô phỏng đã hiệu chỉnh",
        "table": [
            ["Thiết lập", "EA", "PA"],
            ["SFT Vi-only", "68,7", "60,8"],
            ["SFT Vi+FinQA", "70,8", "63,9"],
            ["KD 27B->4B", "72,8", "68,6"],
            ["KD+GRPO/PCPO", "73,6", "70,2"],
            ["KD+GRPO+vote N=15", "74,4", "70,8"],
        ],
        "bullets": [
            "Các dòng này là mô phỏng bảo thủ, chưa phải kết quả chính thức.",
            "Vùng kỳ vọng hợp lý cho pipeline hiện tại: EA 0,73-0,75 và PA 0,70-0,71.",
            "Mốc paper chỉ dùng làm tham chiếu, không dùng để nâng giả định của hệ thống hiện tại.",
        ],
    },
    {
        "title": "15. Phân tích lỗi",
        "bullets": [
            "Nhầm số liệu: sai năm, sai dòng chỉ tiêu hoặc bỏ qua đơn vị.",
            "Sai DSL: đảo thứ tự divide/subtract, tham chiếu # sai, sinh hàm ngoài danh sách.",
            "Đúng EA nhưng sai PA vì chương trình khác gold hoặc PA chưa phủ hết equivalence.",
            "Reasoning hallucination: giải thích nhắc số liệu không có trong input.",
        ],
    },
    {
        "title": "16. Đóng góp cụ thể",
        "bullets": [
            "Pipeline end-to-end cho VLSP 2025 từ dữ liệu đến đánh giá.",
            "Dữ liệu SFT/GRPO/teacher rõ ràng, hỗ trợ program_re và FinQA English.",
            "Distillation reasoning trace giúp student nhỏ học logic tài chính.",
            "PCPO reward ưu tiên program validity, execution và conciseness.",
            "Đánh giá EA/PA symbolic, phù hợp yêu cầu kiểm toán tài chính.",
        ],
    },
    {
        "title": "17. Định hướng phát triển",
        "bullets": [
            "PCPO-Evidence: thưởng chương trình đúng và số liệu map được về bảng/văn bản.",
            "Tăng dữ liệu Table & Text, phép hiếm và counterfactual samples.",
            "Grammar-constrained decoding để giảm invalid program.",
            "Agentic verifier nhẹ: subgoal -> evidence -> DSL -> execute -> vote.",
            "Mở rộng PA bằng nhiều chương trình gold và rewrite rules symbolic.",
        ],
    },
    {
        "title": "18. Kết luận",
        "bullets": [
            "KD + SFT + GRPO là hướng phù hợp nhất khi PA và tính kiểm toán là ưu tiên.",
            "Agent decomposition nên được bổ sung cho Table & Text để tăng generalization.",
            "Hệ thống hiện tại đã có nền tảng tốt; cần hoàn thiện reward, dữ liệu và verifier.",
            "Kết quả chính thức nên được tái tạo bằng pipeline đầy đủ trên RTX 6000/A100.",
        ],
    },
]


def split_table_row(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def is_table_sep(line: str) -> bool:
    stripped = line.strip().strip("|").strip()
    if not stripped:
        return False
    return all(set(cell.strip()) <= {"-", ":"} for cell in stripped.split("|"))


def add_markdown_to_docx(document: Document, markdown: str) -> None:
    lines = markdown.splitlines()
    index = 0
    in_code = False
    code_lines: list[str] = []
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                paragraph = document.add_paragraph()
                run = paragraph.add_run("\n".join(code_lines))
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                paragraph.paragraph_format.left_indent = Cm(0.5)
                code_lines = []
                in_code = False
            else:
                in_code = True
            index += 1
            continue
        if in_code:
            code_lines.append(line)
            index += 1
            continue

        if not stripped:
            index += 1
            continue

        if stripped.startswith("|") and "|" in stripped[1:]:
            rows = []
            while index < len(lines) and lines[index].strip().startswith("|"):
                if not is_table_sep(lines[index]):
                    rows.append(split_table_row(lines[index]))
                index += 1
            if rows:
                max_cols = max(len(row) for row in rows)
                table = document.add_table(rows=len(rows), cols=max_cols)
                table.style = "Table Grid"
                for r_idx, row in enumerate(rows):
                    for c_idx in range(max_cols):
                        text = row[c_idx] if c_idx < len(row) else ""
                        cell = table.cell(r_idx, c_idx)
                        cell.text = text
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = "Times New Roman"
                                run.font.size = Pt(9)
                                if r_idx == 0:
                                    run.bold = True
                    if r_idx == 0:
                        for cell in table.rows[r_idx].cells:
                            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                document.add_paragraph()
            continue

        if stripped.startswith("# "):
            paragraph = document.add_paragraph(stripped[2:].strip(), style="Title")
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            index += 1
            continue
        if stripped.startswith("## "):
            document.add_paragraph(stripped[3:].strip(), style="Heading 1")
            index += 1
            continue
        if stripped.startswith("### "):
            document.add_paragraph(stripped[4:].strip(), style="Heading 2")
            index += 1
            continue
        if stripped.startswith("#### "):
            document.add_paragraph(stripped[5:].strip(), style="Heading 3")
            index += 1
            continue

        if stripped.startswith("- "):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
            index += 1
            continue
        if len(stripped) > 3 and stripped[0].isdigit() and ". " in stripped[:5]:
            document.add_paragraph(stripped.split(". ", 1)[1].strip(), style="List Number")
            index += 1
            continue

        paragraph = document.add_paragraph()
        run = paragraph.add_run(stripped)
        run.font.name = "Times New Roman"
        run.font.size = Pt(11)
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.line_spacing = 1.15
        index += 1


def build_docx() -> None:
    document = Document()
    section = document.sections[0]
    section.top_margin = Cm(2.0)
    section.bottom_margin = Cm(2.0)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.0)

    styles = document.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(11)
    for style_name, size, color in [
        ("Title", 18, RGBColor(160, 0, 28)),
        ("Heading 1", 15, RGBColor(160, 0, 28)),
        ("Heading 2", 13, RGBColor(120, 0, 20)),
        ("Heading 3", 12, RGBColor(80, 80, 80)),
    ]:
        style = styles[style_name]
        style.font.name = "Times New Roman"
        style.font.size = Pt(size)
        style.font.color.rgb = color
        style.font.bold = True

    add_markdown_to_docx(document, REPORT_MD)

    for section in document.sections:
        footer = section.footer.paragraphs[0]
        footer.text = "VLSP 2025 NumQA - Knowledge Distillation Research Report"
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in footer.runs:
            run.font.size = Pt(8)
            run.font.name = "Times New Roman"
            run.font.color.rgb = RGBColor(100, 100, 100)

    document.save(DOCX_PATH)


def remove_all_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public remove API
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=(0, 0, 0), align=None):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = PptRGBColor(*color)
    return shape


def add_bullets(slide, bullets: Iterable[str], left, top, width, height, size=18):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(size)
        paragraph.font.name = "Arial"
        paragraph.space_after = PptPt(6)
    return shape


def add_table(slide, rows: list[list[str]], left, top, width, height, font_size=12):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            text_frame = cell.text_frame
            text_frame.margin_left = PptInches(0.03)
            text_frame.margin_right = PptInches(0.03)
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = PptPt(font_size)
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = PptRGBColor(255, 255, 255)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptRGBColor(165, 0, 33)
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptRGBColor(245, 245, 245)
    return table_shape


def decorate_slide(slide, title: str, number: int):
    width = PptInches(10)
    slide.shapes.add_shape(1, 0, 0, width, PptInches(0.22)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = PptRGBColor(165, 0, 33)
    add_textbox(slide, PptInches(0.35), PptInches(0.32), PptInches(8.8), PptInches(0.55), title, size=22, bold=True, color=(165, 0, 33))
    add_textbox(slide, PptInches(0.35), PptInches(7.12), PptInches(7.7), PptInches(0.22), "Đại học Bách khoa Hà Nội | VLSP 2025 NumQA", size=8, color=(95, 95, 95))
    add_textbox(slide, PptInches(9.15), PptInches(7.08), PptInches(0.5), PptInches(0.25), str(number), size=9, bold=True, color=(165, 0, 33), align=PP_ALIGN.RIGHT)


def build_pptx() -> None:
    if TEMPLATE_PPTX.exists():
        prs = Presentation(str(TEMPLATE_PPTX))
        remove_all_slides(prs)
        blank_layout = prs.slide_layouts[8] if len(prs.slide_layouts) > 8 else prs.slide_layouts[0]
    else:
        prs = Presentation()
        prs.slide_width = PptInches(10)
        prs.slide_height = PptInches(7.5)
        blank_layout = prs.slide_layouts[6]

    for index, spec in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank_layout)
        if index == 1:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = PptRGBColor(255, 255, 255)
            slide.shapes.add_shape(1, 0, 0, PptInches(10), PptInches(0.42)).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = PptRGBColor(165, 0, 33)
            add_textbox(slide, PptInches(0.55), PptInches(1.05), PptInches(8.9), PptInches(1.0), spec["title"], size=30, bold=True, color=(165, 0, 33), align=PP_ALIGN.CENTER)
            add_textbox(slide, PptInches(1.0), PptInches(2.15), PptInches(8.0), PptInches(0.85), spec["subtitle"], size=17, color=(50, 50, 50), align=PP_ALIGN.CENTER)
            add_bullets(slide, spec["bullets"], PptInches(2.0), PptInches(3.55), PptInches(6.5), PptInches(1.55), size=16)
            add_textbox(slide, PptInches(0.7), PptInches(6.7), PptInches(8.8), PptInches(0.35), "HUST template | 22/05/2026", size=10, color=(95, 95, 95), align=PP_ALIGN.CENTER)
            continue

        decorate_slide(slide, spec["title"], index)
        if "table" in spec:
            top = PptInches(1.15)
            table_height = PptInches(2.6 if len(spec["table"]) <= 5 else 3.1)
            add_table(slide, spec["table"], PptInches(0.55), top, PptInches(8.9), table_height, font_size=10 if len(spec["table"]) > 5 else 12)
            if spec.get("bullets"):
                add_bullets(slide, spec["bullets"], PptInches(0.75), top + table_height + PptInches(0.35), PptInches(8.5), PptInches(1.6), size=15)
        else:
            add_bullets(slide, spec.get("bullets", []), PptInches(0.75), PptInches(1.25), PptInches(8.6), PptInches(5.3), size=17)

    prs.save(PPTX_PATH)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    MD_PATH.write_text(REPORT_MD + "\n", encoding="utf-8")
    build_docx()
    build_pptx()
    for source, target in [
        (MD_PATH, CANONICAL_MD_PATH),
        (DOCX_PATH, CANONICAL_DOCX_PATH),
        (PPTX_PATH, CANONICAL_PPTX_PATH),
    ]:
        try:
            shutil.copy2(source, target)
            print(f"updated canonical: {target}")
        except PermissionError:
            print(f"canonical locked, kept revised copy: {target}")
    print(MD_PATH)
    print(DOCX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()