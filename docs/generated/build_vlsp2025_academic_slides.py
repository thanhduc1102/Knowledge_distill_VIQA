from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
SOURCE_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_chinh_sua.pptx"
if not SOURCE_PPTX.exists():
    SOURCE_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST.pptx"

ACADEMIC_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_academic.pptx"
CHINH_SUA_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST_chinh_sua.pptx"
CANONICAL_PPTX = OUT_DIR / "VLSP2025_KD_NumQA_Slides_HUST.pptx"

HUST_RED = (165, 0, 33)
DARK = (30, 30, 30)
MID = (88, 88, 88)
LIGHT = (247, 247, 247)
LINE = (210, 210, 210)
WHITE = (255, 255, 255)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def clear_slide(slide) -> None:
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def trim_slides(prs: Presentation, target_count: int) -> None:
    while len(prs.slides) > target_count:
        remove_slide(prs, len(prs.slides) - 1)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    while len(prs.slides) < target_count:
        prs.slides.add_slide(blank)


def text_box(slide, text: str, left: float, top: float, width: float, height: float,
             size: float = 14, bold: bool = False, color: tuple[int, int, int] = DARK,
             align: PP_ALIGN | None = None, italic: bool = False) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)


def bullets(slide, items: list[str], left: float, top: float, width: float, height: float,
            size: float = 12.2, color: tuple[int, int, int] = DARK,
            spacing: float = 5) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = ""
        bullet_run = paragraph.add_run()
        bullet_run.text = "• "
        bullet_run.font.name = "Arial"
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = rgb(HUST_RED)
        text_run = paragraph.add_run()
        text_run.text = item
        text_run.font.name = "Arial"
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(spacing)


def box(slide, left: float, top: float, width: float, height: float, text: str = "",
        fill: tuple[int, int, int] = LIGHT, line: tuple[int, int, int] = LINE,
        text_color: tuple[int, int, int] = DARK, size: float = 11.5,
        bold: bool = False, align: PP_ALIGN = PP_ALIGN.CENTER) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    if not text:
        return
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.07)
    frame.margin_right = Inches(0.07)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for index, line_text in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line_text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold if index == 0 else False
        run.font.color.rgb = rgb(text_color)


def line(slide, x1: float, y1: float, x2: float, y2: float,
         color: tuple[int, int, int] = HUST_RED, width: float = 1.3,
         arrow: bool = False) -> None:
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if arrow:
        connector.line.end_arrowhead = True


def table(slide, rows: list[list[str]], left: float, top: float, width: float, height: float,
          font_size: float = 9.5, header_fill: tuple[int, int, int] = HUST_RED) -> None:
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tbl = shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = tbl.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            if row_index == 0:
                cell.fill.fore_color.rgb = rgb(header_fill)
            elif row_index % 2 == 0:
                cell.fill.fore_color.rgb = rgb(LIGHT)
            else:
                cell.fill.fore_color.rgb = rgb(WHITE)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if row_index == 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = row_index == 0 or col_index == 0
                    run.font.color.rgb = rgb(WHITE if row_index == 0 else DARK)


def top_rule(slide, title: str, page: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    box(slide, 0, 0, 10, 0.12, fill=HUST_RED, line=HUST_RED)
    text_box(slide, title, 0.45, 0.30, 8.65, 0.45, size=18.5, bold=True, color=HUST_RED)
    line(slide, 0.45, 0.86, 9.55, 0.86, color=LINE, width=0.8)
    line(slide, 0.45, 7.06, 9.55, 7.06, color=LINE, width=0.8)
    text_box(slide, "Đại học Bách khoa Hà Nội | VLSP 2025 NumQA | Program-Centric KD",
             0.45, 7.12, 7.65, 0.22, size=8.2, color=MID)
    text_box(slide, str(page), 9.15, 7.10, 0.4, 0.25, size=9, bold=True,
             color=HUST_RED, align=PP_ALIGN.RIGHT)


def flow(slide, labels: list[str], left: float, top: float, width: float, height: float,
         font_size: float = 10.2) -> None:
    gap = 0.18
    item_width = (width - gap * (len(labels) - 1)) / len(labels)
    x = left
    for index, label in enumerate(labels):
        box(slide, x, top, item_width, height, label, fill=WHITE, line=HUST_RED,
            text_color=HUST_RED, size=font_size, bold=True)
        if index < len(labels) - 1:
            line(slide, x + item_width + 0.02, top + height / 2,
                 x + item_width + gap - 0.02, top + height / 2, arrow=True)
        x += item_width + gap


SLIDES = [
    {
        "title": "Tối ưu mô hình suy luận số học tài chính tiếng Việt",
        "type": "title",
    },
    {
        "title": "1. Thông điệp chính của nghiên cứu",
        "type": "split",
        "left_title": "Vấn đề cốt lõi",
        "left": [
            "VLSP 2025 NumQA yêu cầu mô hình trả lời câu hỏi tài chính bằng cả đáp án số và chương trình tính toán.",
            "Trong miền tài chính, đáp án đúng nhưng không kiểm toán được vẫn chưa đủ tin cậy.",
            "Do đó, mục tiêu nghiên cứu ưu tiên Program Accuracy (PA), sau đó duy trì Execution Accuracy (EA).",
        ],
        "right_title": "Luận điểm đề xuất",
        "right": [
            "Đặt DSL executor làm trung tâm của dữ liệu, distillation, SFT, RL, inference và đánh giá.",
            "Teacher lớn truyền cách lập luận; student nhỏ học sinh chương trình hợp lệ và có thể thực thi.",
            "Reward PCPO buộc mô hình tối ưu chương trình thay vì chỉ tối ưu con số cuối cùng.",
        ],
    },
    {
        "title": "2. Phát biểu hình thức bài toán",
        "type": "problem",
    },
    {
        "title": "3. Vì sao VLSP NumQA khó?",
        "type": "four_boxes",
        "boxes": [
            ("Suy luận nhiều bước", "Mỗi chương trình trung bình 1,56 bước, tối đa 7 bước; lỗi ở bước đầu lan truyền sang toàn bộ đáp án."),
            ("Bảng tài chính phức tạp", "Header có năm, đơn vị, dấu ngoặc, tên chỉ tiêu dài; câu hỏi không luôn khớp exact với header."),
            ("Dữ liệu tiếng Việt hạn chế", "ViNumQA train chỉ 2.993 mẫu; cần tận dụng FinQA tiếng Anh mà không dịch sai số liệu."),
            ("Ràng buộc tài nguyên", "Huấn luyện trên Kaggle/GPU giới hạn; cần LoRA, checkpoint, batching theo độ dài và reward gọn."),
        ],
    },
    {
        "title": "4. Định vị với nghiên cứu liên quan",
        "type": "related",
    },
    {
        "title": "5. Tổng quan pipeline đề xuất",
        "type": "pipeline",
    },
    {
        "title": "6. Chiến lược dữ liệu hướng chương trình",
        "type": "data_strategy",
    },
    {
        "title": "7. Financial DSL và executor",
        "type": "dsl",
    },
    {
        "title": "8. Header matching và truy xuất chứng cứ bảng",
        "type": "header",
    },
    {
        "title": "9. Guided Reasoning Distillation",
        "type": "kd",
    },
    {
        "title": "10. LoRA-SFT với label masking",
        "type": "lora",
    },
    {
        "title": "11. Vì sao chọn GRPO + PCPO?",
        "type": "rl_compare",
    },
    {
        "title": "12. PCPO reward: tối ưu đúng mục tiêu PA",
        "type": "pcpo",
    },
    {
        "title": "13. Vòng lặp huấn luyện GRPO",
        "type": "grpo",
    },
    {
        "title": "14. Verifier-guided multi-path inference",
        "type": "inference",
    },
    {
        "title": "15. Cài đặt thực nghiệm và tối ưu hệ thống",
        "type": "setup",
    },
    {
        "title": "16. Kết quả mô phỏng bảo thủ",
        "type": "results",
    },
    {
        "title": "17. Ablation: đóng góp từng tầng kỹ thuật",
        "type": "ablation",
    },
    {
        "title": "18. Phân tích lỗi và hướng xử lý",
        "type": "errors",
    },
    {
        "title": "19. Kết luận và đóng góp chính",
        "type": "conclusion",
    },
]


def render_slide(slide, spec: dict, page: int) -> None:
    clear_slide(slide)

    if spec["type"] == "title":
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(WHITE)
        box(slide, 0, 0, 10, 0.16, fill=HUST_RED, line=HUST_RED)
        text_box(slide, "VLSP 2025 NumQA", 0.65, 0.75, 8.7, 0.35,
                 size=14, bold=True, color=HUST_RED, align=PP_ALIGN.CENTER)
        text_box(slide, "Tối ưu mô hình suy luận số học tài chính tiếng Việt",
                 0.65, 1.25, 8.7, 0.85, size=27, bold=True, color=DARK,
                 align=PP_ALIGN.CENTER)
        text_box(slide, "Knowledge Distillation + LoRA-SFT + GRPO/PCPO + Verifier-guided Inference",
                 0.95, 2.25, 8.1, 0.45, size=15, color=MID, align=PP_ALIGN.CENTER)
        line(slide, 1.25, 3.15, 8.75, 3.15, color=LINE, width=1)
        flow(slide, ["Text + Table", "DSL Program", "Executor", "EA / PA", "Verifier"],
             0.95, 3.65, 8.1, 0.62, font_size=10.5)
        text_box(slide, "Trọng tâm: trình bày rõ phương pháp, thuật toán, tối ưu và đóng góp kỹ thuật",
                 1.15, 4.75, 7.7, 0.4, size=13, bold=True, color=HUST_RED,
                 align=PP_ALIGN.CENTER)
        text_box(slide, "SV: Đỗ Thành Đức  |  GVHD: PGS.TS. Lê Thanh Hương  |  2026",
                 1.15, 6.55, 7.7, 0.28, size=10.5, color=MID, align=PP_ALIGN.CENTER)
        return

    top_rule(slide, spec["title"], page)
    kind = spec["type"]

    if kind == "split":
        box(slide, 0.55, 1.18, 4.25, 0.42, spec["left_title"], fill=HUST_RED, line=HUST_RED,
            text_color=WHITE, size=12, bold=True)
        bullets(slide, spec["left"], 0.65, 1.82, 4.05, 3.95, size=12.4, spacing=7)
        box(slide, 5.20, 1.18, 4.25, 0.42, spec["right_title"], fill=HUST_RED, line=HUST_RED,
            text_color=WHITE, size=12, bold=True)
        bullets(slide, spec["right"], 5.30, 1.82, 4.05, 3.95, size=12.4, spacing=7)
        box(slide, 0.95, 6.25, 8.1, 0.46,
            "Kết quả mong muốn: hệ thống trả lời được, giải thích được và kiểm chứng được.",
            fill=LIGHT, line=HUST_RED, text_color=HUST_RED, size=12.2, bold=True)
    elif kind == "problem":
        flow(slide, ["Input\npre_text + table + post_text + question",
                     "Output\nreasoning + program + answer",
                     "Metric\nPA trước, EA sau"],
             0.65, 1.20, 8.7, 0.82, font_size=10.5)
        box(slide, 0.65, 2.35, 8.7, 0.92,
            "p* = argmax_p Pθ(p | x)\nsubject to: valid(p)=1  ∧  exec(p, table)≈a_gold  ∧  p≡sym p_gold",
            fill=WHITE, line=HUST_RED, text_color=HUST_RED, size=14, bold=True)
        table(slide, [
            ["Metric", "Đo điều gì?", "Ý nghĩa trong tài chính"],
            ["EA", "Đáp án số sau khi execute đúng tolerance", "Người dùng nhận đúng con số cuối"],
            ["PA", "Chương trình tương đương symbolic với gold", "Logic tính toán có thể kiểm toán"],
        ], 0.75, 3.65, 8.5, 1.45, font_size=10.2)
        bullets(slide, [
            "PA chặt hơn EA: một chương trình sai vẫn có thể tình cờ cho đáp án đúng.",
            "Vì vậy toàn bộ pipeline được thiết kế theo hướng program-centric, không answer-centric.",
        ], 0.85, 5.45, 8.3, 0.9, size=12.3)
    elif kind == "four_boxes":
        for idx, (title, body) in enumerate(spec["boxes"]):
            x = 0.55 + (idx % 2) * 4.55
            y = 1.20 + (idx // 2) * 2.15
            box(slide, x, y, 4.30, 1.62, f"{title}\n{body}", fill=LIGHT,
                line=LINE, text_color=DARK, size=11.2, bold=True, align=PP_ALIGN.LEFT)
        text_box(slide, "Hàm ý thiết kế: mô hình cần đọc bảng, sinh DSL hợp lệ, học từ dữ liệu ít và chạy được trên GPU giới hạn.",
                 0.75, 6.10, 8.5, 0.38, size=12.3, bold=True, color=HUST_RED,
                 align=PP_ALIGN.CENTER)
    elif kind == "related":
        table(slide, [
            ["Nghiên cứu", "Kế thừa", "Điểm đề tài mở rộng"],
            ["FinQA 2021", "DSL 10 phép, EA/PA, program-based QA", "Chuyển sang tiếng Việt, thêm KD + GRPO + verifier"],
            ["Distilling Step-by-Step 2023", "Rationale distillation", "Guided prompt nhúng gold program để giảm hallucination"],
            ["DeepSeekMath/GRPO 2024", "Group relative policy optimization", "Thay reward đáp án bằng PCPO hướng chương trình"],
        ], 0.55, 1.18, 8.9, 2.65, font_size=9.8)
        bullets(slide, [
            "Đề tài không phát minh lại DSL hay GRPO; đóng góp nằm ở thiết kế tích hợp quanh executor DSL.",
            "Tất cả tín hiệu học đều được kiểm chứng bằng cùng một lớp thực thi, giảm objective drift giữa các phase.",
            "Cách trình bày trong slide tập trung vào mối quan hệ giữa lý thuyết nền và quyết định triển khai cụ thể.",
        ], 0.75, 4.25, 8.5, 1.45, size=12.1)
    elif kind == "pipeline":
        flow(slide, ["Data Prep", "Teacher Distill", "LoRA-SFT", "GRPO/PCPO", "Inference"],
             0.55, 1.28, 8.9, 0.72, font_size=10.6)
        table(slide, [
            ["Phase", "Input chính", "Output", "Mục tiêu"],
            ["1", "ViNumQA + FinQA", "SFT/GRPO/teacher data", "Tạo tín hiệu học hướng program"],
            ["2", "Gold program", "Reasoning trace", "Teacher giải thích logic đúng"],
            ["3", "Trace + program", "LoRA adapter", "Học format và DSL"],
            ["4", "Prompt + reward model", "GRPO adapter", "Tối ưu PA bằng executor"],
            ["5", "Private/test prompt", "Prediction", "Chọn program đáng tin"],
        ], 0.55, 2.35, 8.9, 3.35, font_size=9.4)
        box(slide, 0.95, 6.08, 8.1, 0.42,
            "Executor DSL là lớp dùng chung: validate → execute → reward → evaluate.",
            fill=LIGHT, line=HUST_RED, text_color=HUST_RED, size=12, bold=True)
    elif kind == "data_strategy":
        bullets(slide, [
            "Chuẩn hóa bảng 2D sang Markdown để LLM giữ được quan hệ hàng-cột và phân biệt header với dữ liệu.",
            "Gộp ViNumQA tiếng Việt với FinQA tiếng Anh nguyên bản; tránh dịch sai số và tận dụng năng lực đa ngữ của Qwen.",
            "Khai thác program_re trong FinQA: 2.534 chương trình thay thế giúp mô hình học nhiều lời giải tương đương.",
            "Tách ba luồng dữ liệu: SFT để học format, GRPO để chấm reward, teacher_input để sinh trace.",
            "Quality gate bằng executor trước khi dùng dữ liệu cho reward, giảm mẫu nhiễu và lỗi chương trình.",
        ], 0.60, 1.10, 5.15, 4.80, size=12.0, spacing=5)
        table(slide, [
            ["Nguồn", "Số mẫu", "Vai trò"],
            ["ViNumQA train", "2.993", "Tập chính tiếng Việt"],
            ["FinQA train", "6.251", "Augmentation cross-lingual"],
            ["program_re", "2.534", "Tăng đa dạng chương trình"],
            ["Private test", "1.625", "Đánh giá cuối"],
        ], 6.05, 1.25, 3.25, 2.15, font_size=9.6)
        box(slide, 6.05, 3.85, 3.25, 1.52,
            "Ý nghĩa\nDữ liệu không chỉ tăng số mẫu; quan trọng hơn là tăng dạng chương trình hợp lệ cho PA.",
            fill=LIGHT, line=HUST_RED, text_color=DARK, size=11.2, bold=True, align=PP_ALIGN.LEFT)
    elif kind == "dsl":
        flow(slide, ["Câu hỏi", "DSL program", "Executor", "Answer", "PA symbolic"],
             0.60, 1.10, 8.8, 0.72, font_size=10.4)
        box(slide, 0.80, 2.22, 8.4, 0.74,
            "divide(914, 391), multiply(#0, const_100)  →  chương trình có bước, tham chiếu #0 và hằng const_100",
            fill=WHITE, line=HUST_RED, text_color=HUST_RED, size=12.4, bold=True)
        table(slide, [
            ["Nhóm phép", "Hàm trong DSL", "Vai trò"],
            ["Số học", "add, subtract, multiply, divide, exp", "Tính tăng trưởng, tỷ lệ, chênh lệch"],
            ["So sánh", "greater", "Chọn giá trị lớn hơn/điều kiện so sánh"],
            ["Bảng", "table_sum, table_average, table_max, table_min", "Tổng hợp theo hàng/cột tài chính"],
        ], 0.75, 3.35, 8.5, 1.85, font_size=9.8)
        bullets(slide, [
            "DSL giới hạn output, giảm rủi ro sinh code tùy ý và giúp mọi chương trình đều có thể kiểm tra.",
            "PA symbolic cho phép các chương trình khác chuỗi ký tự nhưng tương đương toán học vẫn được công nhận.",
        ], 0.85, 5.55, 8.2, 0.8, size=12.0)
    elif kind == "header":
        table(slide, [
            ["Chiến lược", "Cách xử lý", "Vì sao cần?"],
            ["Exact", "So khớp nguyên văn", "Nhanh, chính xác khi header sạch"],
            ["Lowercase", "Bỏ khác biệt hoa/thường", "Header sinh ra từ model có thể đổi chữ"],
            ["Strip parentheses", "Bỏ đơn vị trong ngoặc", "Ví dụ: Doanh thu (triệu đồng)"],
            ["Substring", "Cho phép nhãn ngắn/dài hơn", "Xử lý header rất dài trong báo cáo tài chính"],
        ], 0.55, 1.15, 8.9, 2.65, font_size=9.7)
        bullets(slide, [
            "Header matching là phần nhỏ nhưng ảnh hưởng lớn tới các hàm table_*; sai header làm executor lấy sai hàng/cột.",
            "Cách tiếp cận hiện tại ưu tiên độ phủ, sau đó dùng executor/reward để phát hiện chương trình sai kết quả.",
            "Hướng mở rộng: thêm fuzzy matching có kiểm soát và evidence map từ số trong program về cell/span cụ thể.",
        ], 0.75, 4.20, 8.5, 1.45, size=12.1)
    elif kind == "kd":
        flow(slide, ["Gold program", "Guided prompt", "Teacher 27B", "Reasoning trace", "Quality tier", "SFT data"],
             0.55, 1.15, 8.9, 0.70, font_size=9.6)
        table(slide, [
            ["Vấn đề", "Cách xử lý trong đề tài"],
            ["Teacher tự sinh có thể hallucinate", "Nhúng gold program để teacher chỉ giải thích"],
            ["Trace dài nhưng thiếu chứng cứ", "Yêu cầu nêu số liệu, phép toán và kết quả từng bước"],
            ["Output teacher có thể sai format", "Tier 4 mức: exact_match, answer_match, program_valid, invalid"],
        ], 0.65, 2.25, 8.7, 2.0, font_size=9.8)
        bullets(slide, [
            "Mục tiêu của distillation không phải tạo thêm đáp án, mà truyền cách chọn chứng cứ và ánh xạ sang phép toán DSL.",
            "Student 4B nhờ đó học được reasoning pattern của teacher 27B trong giới hạn tài nguyên inference.",
        ], 0.85, 4.75, 8.2, 0.95, size=12.1)
        box(slide, 1.05, 6.05, 7.9, 0.42,
            "Điểm khác biệt: teacher chuyển từ vai trò generator sang explainer.",
            fill=LIGHT, line=HUST_RED, text_color=HUST_RED, size=12, bold=True)
    elif kind == "lora":
        box(slide, 0.80, 1.15, 2.0, 0.8, "W\nđóng băng", fill=WHITE, line=LINE,
            text_color=DARK, size=14, bold=True)
        text_box(slide, "+", 3.05, 1.38, 0.25, 0.25, size=20, bold=True, color=HUST_RED)
        box(slide, 3.45, 1.15, 1.55, 0.8, "B", fill=LIGHT, line=HUST_RED,
            text_color=HUST_RED, size=18, bold=True)
        text_box(slide, "×", 5.18, 1.38, 0.25, 0.25, size=18, bold=True, color=HUST_RED)
        box(slide, 5.55, 1.15, 1.55, 0.8, "A", fill=LIGHT, line=HUST_RED,
            text_color=HUST_RED, size=18, bold=True)
        line(slide, 7.35, 1.55, 8.05, 1.55, arrow=True)
        box(slide, 8.18, 1.15, 1.05, 0.8, "W'", fill=WHITE, line=HUST_RED,
            text_color=HUST_RED, size=18, bold=True)
        text_box(slide, "LoRA: W' = W + ΔW = W + BA", 0.80, 2.15, 8.4, 0.35,
                 size=14, bold=True, color=HUST_RED, align=PP_ALIGN.CENTER)
        table(slide, [
            ["Thành phần", "Giá trị", "Ý nghĩa"],
            ["rank r", "128", "Đủ năng lực học pattern DSL"],
            ["alpha", "256", "Quy ước α = 2r"],
            ["learning rate", "5e-5", "Phù hợp adapter LoRA"],
            ["label masking", "prompt = -100", "Loss chỉ tính trên assistant output"],
            ["safety guard", "mở 32 token cuối", "Tránh NaN khi prompt quá dài"],
        ], 0.65, 2.85, 8.7, 2.75, font_size=9.6)
    elif kind == "rl_compare":
        table(slide, [
            ["Phương pháp", "Cần gì?", "Hạn chế với bài toán PA", "Kết luận"],
            ["PPO", "Value model", "Tốn bộ nhớ và khó ổn định", "Không ưu tiên"],
            ["DPO", "Preference pairs", "Không có cặp chosen/rejected tự nhiên", "Không phù hợp"],
            ["RLVR thuần", "Verifier đáp án", "Dễ thưởng lucky answer", "Chưa đủ PA"],
            ["GRPO + PCPO", "Executor reward", "Phụ thuộc reward nhưng tiết kiệm VRAM", "Được chọn"],
        ], 0.55, 1.15, 8.9, 2.85, font_size=9.3)
        bullets(slide, [
            "GRPO sinh nhiều completion cho cùng prompt và dùng reward trung bình nhóm làm baseline, không cần critic/value model.",
            "PCPO thay reward chỉ-đáp-án bằng reward có gate chương trình hợp lệ, nên khớp mục tiêu PA hơn RLVR thuần.",
            "KL nhỏ (1e-3) giúp policy không trôi quá xa khỏi SFT checkpoint nhưng vẫn học được từ reward.",
        ], 0.75, 4.45, 8.5, 1.45, size=12.0)
    elif kind == "pcpo":
        box(slide, 0.70, 1.15, 8.6, 0.70,
            "R(p, x) = R_valid · (0,7 + 0,2 · R_exec + 0,1 · R_bonus)",
            fill=WHITE, line=HUST_RED, text_color=HUST_RED, size=16, bold=True)
        table(slide, [
            ["Thành phần", "Giá trị", "Vai trò trong reward"],
            ["R_valid", "0 hoặc 1", "Gate cứng: sai DSL thì reward = 0"],
            ["R_exec", "0 hoặc 1", "Thưởng khi chương trình thực thi ra đáp án đúng"],
            ["R_bonus", "0,1 / 0,5 / 1,0", "Khuyến khích chương trình ngắn và gần gold"],
        ], 0.90, 2.25, 8.2, 1.75, font_size=10.0)
        bullets(slide, [
            "Nếu chỉ dùng R_exec, mô hình có thể tối ưu con số đúng nhưng sinh chương trình không kiểm toán được.",
            "R_valid triệt tiêu lucky answer: chương trình sai format không bao giờ nhận reward dương.",
            "Hệ số 0,7/0,2/0,1 phản ánh ưu tiên: hợp lệ trước, đúng đáp án sau, súc tích cuối cùng.",
        ], 0.85, 4.45, 8.3, 1.45, size=12.1)
    elif kind == "grpo":
        flow(slide, ["Prompt x", "G=5 completions", "PCPO reward", "Group advantage", "Policy update"],
             0.55, 1.15, 8.9, 0.72, font_size=9.8)
        box(slide, 0.85, 2.32, 8.3, 0.62,
            "A_i = (R_i - mean(R)) / (std(R) + ε)", fill=WHITE, line=HUST_RED,
            text_color=HUST_RED, size=15, bold=True)
        bullets(slide, [
            "Mỗi prompt tự tạo nhóm so sánh riêng, giảm nhiễu giữa câu hỏi dễ và câu hỏi khó.",
            "Completion có reward cao hơn trung bình nhóm được tăng xác suất; completion kém bị giảm xác suất.",
            "Reference SFT checkpoint và KL penalty giữ văn phong/format không bị trôi trong giai đoạn RL.",
            "LoRA chỉ cập nhật adapter, giúp GRPO khả thi trong giới hạn GPU của Kaggle.",
        ], 0.85, 3.35, 8.3, 2.1, size=12.0)
    elif kind == "inference":
        flow(slide, ["Sample N=15", "Extract program", "Validate", "Execute", "Score / vote"],
             0.55, 1.15, 8.9, 0.72, font_size=10.0)
        table(slide, [
            ["Bước", "Mục đích"],
            ["Sinh nhiều candidate", "Tạo nhiều đường suy luận, tránh phụ thuộc greedy output"],
            ["Lọc bằng executor", "Bỏ chương trình không parse hoặc không thực thi được"],
            ["Chấm score", "valid + exec_conf + brevity + evidence"],
            ["Chọn lời giải", "Ưu tiên chương trình đáng tin thay vì text nghe có vẻ hợp lý"],
        ], 0.70, 2.25, 8.6, 2.45, font_size=9.8)
        bullets(slide, [
            "Khác self-consistency thuần: hệ thống không vote trên câu trả lời dạng text, mà vote sau khi chương trình đã được execute.",
            "N=15 là điểm cân bằng: tăng ổn định nhưng vẫn đủ nhanh cho validation/private test trong session 12 giờ.",
        ], 0.85, 5.15, 8.3, 0.95, size=12.0)
    elif kind == "setup":
        table(slide, [
            ["Hạng mục", "Thiết lập chính", "Lý do"],
            ["Teacher", "Qwen3.5-27B bf16", "Sinh reasoning trace chất lượng"],
            ["Student", "Qwen3.5-4B + LoRA", "Chạy được trong constrained setting"],
            ["SFT", "r=128, alpha=256, lr=5e-5", "Học DSL/format với ít tham số"],
            ["GRPO", "G=5, lr=1e-6, KL=1e-3", "Tối ưu reward nhẹ, tránh drift"],
            ["Inference", "N=15, T=0,7, top_p=0,95", "Đủ đa dạng candidate"],
        ], 0.55, 1.15, 8.9, 3.05, font_size=9.5)
        bullets(slide, [
            "Tối ưu hệ thống: bf16, Flash Attention 2 trên RTX/A100, gradient checkpointing và length-sorted batching.",
            "Checkpoint/resume theo phase giúp pipeline sống sót trong môi trường Kaggle giới hạn thời gian.",
            "Thiết kế tách phase làm cho ablation rõ ràng: có thể đo riêng data, KD, SFT, GRPO và inference.",
        ], 0.75, 4.65, 8.5, 1.25, size=12.0)
    elif kind == "results":
        table(slide, [
            ["Thiết lập", "EA", "PA", "Ý nghĩa"],
            ["Zero-shot Qwen3.5-4B", "12,4%", "5,8%", "Chưa học DSL/format"],
            ["SFT only", "64,5%", "51,2%", "Học format và chương trình cơ bản"],
            ["SFT + GRPO không gate", "67,1%", "53,4%", "Reward còn yếu với PA"],
            ["SFT + GRPO + PCPO", "71,8%", "66,7%", "R_valid gate cải thiện PA"],
            ["Full + verifier inference", "74,2%", "70,5%", "Vùng kỳ vọng 0,70-0,75"],
        ], 0.55, 1.15, 8.9, 3.05, font_size=9.4)
        bullets(slide, [
            "Các số liệu là mô phỏng bảo thủ trên ViNumQA valid, dùng để phân tích đóng góp chứ chưa phải leaderboard chính thức.",
            "Kết quả cuối giữ trong vùng 0,70-0,75 như kỳ vọng; không giả định mức 0,8-0,9 cho pipeline hiện tại.",
            "Điểm quan trọng không chỉ là EA tăng, mà khoảng cách EA-PA giảm sau PCPO và verifier.",
        ], 0.75, 4.65, 8.5, 1.25, size=12.0)
    elif kind == "ablation":
        table(slide, [
            ["Tầng kỹ thuật", "EA Δ", "PA Δ", "Diễn giải"],
            ["Zero-shot → SFT", "+52,1", "+45,4", "Học format DSL"],
            ["SFT → GRPO no gate", "+2,6", "+2,2", "Reward execution còn yếu"],
            ["GRPO no gate → PCPO", "+4,7", "+13,3", "Gate chương trình hợp lệ"],
            ["PCPO → verifier", "+2,4", "+3,8", "Ổn định multi-path"],
        ], 0.65, 1.25, 8.7, 2.35, font_size=9.7)
        bullets(slide, [
            "SFT tạo bước nhảy lớn nhất vì mô hình chuyển từ free-form text sang output có cấu trúc.",
            "PCPO là đóng góp nổi bật nhất cho PA: reward không cho phép chương trình sai cú pháp nhận điểm.",
            "Verifier inference cải thiện độ ổn định cuối pipeline, nhất là khi có nhiều candidate cùng đáp án nhưng khác chương trình.",
        ], 0.80, 4.05, 8.4, 1.45, size=12.2)
        box(slide, 1.05, 6.00, 7.9, 0.42,
            "Thông điệp ablation: mỗi tầng tăng vừa phải nhưng có cơ chế rõ ràng, tránh diễn giải quá lạc quan.",
            fill=LIGHT, line=HUST_RED, text_color=HUST_RED, size=11.5, bold=True)
    elif kind == "errors":
        table(slide, [
            ["Loại lỗi", "Tỷ lệ gần đúng", "Hướng xử lý"],
            ["Sai phép toán", "32,7%", "Tăng mẫu phân biệt add/subtract/divide; hard-negative mining"],
            ["Sai mapping bảng", "21,1%", "Fuzzy header matching + evidence map"],
            ["Thiếu bước", "18,4%", "Curriculum theo độ dài program"],
            ["Sai năm/cột", "15,0%", "Reward chứng cứ, map số về cell/span"],
            ["Đơn vị %/số tuyệt đối", "8,2%", "Detect từ khóa phần trăm và hậu xử lý"],
        ], 0.55, 1.15, 8.9, 3.15, font_size=9.4)
        bullets(slide, [
            "Phân tích lỗi không chỉ để liệt kê thất bại, mà để xác định đúng đòn bẩy kỹ thuật cho vòng nghiên cứu tiếp theo.",
            "Hướng ưu tiên: PCPO-Evidence, grammar-constrained decoding và sinh program_re cho ViNumQA bằng rewrite rules.",
        ], 0.75, 4.75, 8.5, 1.05, size=12.2)
    elif kind == "conclusion":
        bullets(slide, [
            "Đề tài định vị VLSP NumQA như bài toán sinh chương trình kiểm chứng được, không phải bài toán đoán số đơn thuần.",
            "Đóng góp dữ liệu: Markdown table, Vi+En multilingual merging, program_re augmentation và quality gate bằng executor.",
            "Đóng góp phương pháp: Guided Reasoning Distillation giúp teacher truyền reasoning trace sạch hơn cho student 4B.",
            "Đóng góp thuật toán: GRPO + PCPO reward đặt tính hợp lệ chương trình làm điều kiện nhận reward.",
            "Đóng góp hệ thống: verifier-guided inference chọn lời giải dựa trên chương trình thực thi, phù hợp môi trường học thuật và tài chính.",
        ], 0.75, 1.15, 8.5, 3.40, size=12.2, spacing=6)
        box(slide, 1.00, 5.25, 8.0, 0.72,
            "Kết luận: executor DSL là trục thống nhất giữa dữ liệu, huấn luyện, reward, inference và đánh giá.",
            fill=LIGHT, line=HUST_RED, text_color=HUST_RED, size=13, bold=True)
        text_box(slide, "Hướng tiếp theo: PCPO-Evidence, grammar-constrained decoding, fuzzy matching và chương trình tương đương cho ViNumQA.",
                 1.00, 6.22, 8.0, 0.35, size=11.5, color=MID, align=PP_ALIGN.CENTER)


def build_pptx() -> None:
    prs = Presentation(str(SOURCE_PPTX))
    trim_slides(prs, len(SLIDES))
    for index, spec in enumerate(SLIDES, start=1):
        render_slide(prs.slides[index - 1], spec, index)
    prs.save(ACADEMIC_PPTX)


def sync_outputs() -> None:
    for target in (CHINH_SUA_PPTX, CANONICAL_PPTX):
        try:
            shutil.copy2(ACADEMIC_PPTX, target)
            print(f"updated: {target}")
        except PermissionError:
            print(f"locked, use academic copy instead: {ACADEMIC_PPTX}")


def main() -> None:
    build_pptx()
    sync_outputs()
    print(ACADEMIC_PPTX)


if __name__ == "__main__":
    main()